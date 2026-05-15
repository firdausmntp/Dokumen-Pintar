"""PPTX (PowerPoint) format handler backed by python-pptx."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from dokumen_pintar.errors import HandlerError, UnsupportedFormatError
from dokumen_pintar.handlers.base import (
    FormatHandler,
    HandlerCapability,
    default_registry,
)


def _slide_title(slide: Any) -> str:
    """Return the title text of a slide or empty string."""
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None
    if title_shape is None:
        return ""
    try:
        if title_shape.has_text_frame:
            return title_shape.text_frame.text or ""
    except Exception:
        return ""
    return ""


def _iter_text_frames(slide: Any) -> list[str]:
    """Collect text from all text-frame-bearing shapes on a slide."""
    texts: list[str] = []
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                txt = shape.text_frame.text
                if txt:
                    texts.append(txt)
        except Exception:
            continue
    return texts


def _shape_kind(shape: Any) -> str:
    """Classify shape into 'text'|'picture'|'table'|'other'."""
    try:
        if shape.has_text_frame:
            return "text"
    except Exception:
        pass
    try:
        if shape.has_table:
            return "table"
    except Exception:
        pass
    try:
        # shape_type 13 == PICTURE
        if (
            getattr(shape, "shape_type", None) is not None
            and str(shape.shape_type).upper().find("PICTURE") != -1
        ):
            return "picture"
    except Exception:
        pass
    return "other"


def _shape_to_dict(shape: Any) -> dict[str, Any]:
    kind = _shape_kind(shape)
    base: dict[str, Any] = {
        "kind": kind,
        "name": getattr(shape, "name", ""),
        "text": "",
    }
    try:
        base["shape_id"] = int(getattr(shape, "shape_id", 0) or 0)
    except Exception:
        base["shape_id"] = 0

    # Position / size (EMU -> int) if available.
    for attr in ("left", "top", "width", "height"):
        try:
            val: Any = getattr(shape, attr, None)
            if val is None:
                continue
            if isinstance(val, int):
                base[attr] = val
            else:
                base[attr] = int(val)
        except Exception:
            pass

    if kind == "text":
        try:
            base["text"] = shape.text_frame.text or ""
        except Exception:
            base["text"] = ""
    elif kind == "table":
        try:
            tbl = shape.table
            rows: list[list[str]] = []
            for row in tbl.rows:
                rows.append([cell.text or "" for cell in row.cells])
            base["rows"] = rows
        except Exception:
            base["rows"] = []
    elif kind == "picture":
        try:
            base["image_filename"] = getattr(shape.image, "filename", "") or ""
            base["content_type"] = getattr(shape.image, "content_type", "") or ""
        except Exception:
            pass

    return base


def _parse_slide_index(token: str) -> int:
    try:
        return int(token)
    except (TypeError, ValueError) as exc:
        raise HandlerError(f"invalid slide index: {token!r}") from exc


def _get_slide(prs: Any, index: int) -> Any:
    slides = list(prs.slides)
    if index < 0 or index >= len(slides):
        raise HandlerError(f"slide index out of range: {index} (slide_count={len(slides)})")
    return slides[index]


class PptxHandler:
    """Handler for Microsoft PowerPoint .pptx files."""

    name: str = "pptx"
    extensions: tuple[str, ...] = (".pptx",)
    capabilities: HandlerCapability = (
        HandlerCapability.READ_TEXT
        | HandlerCapability.STRUCTURED_GET
        | HandlerCapability.STRUCTURED_SET
        | HandlerCapability.STRUCTURED_DELETE
        | HandlerCapability.SEARCH_EXTRACTED
        | HandlerCapability.WRITE_META
    )

    _WRITABLE_CORE_PROPS: tuple[str, ...] = (
        "author",
        "category",
        "comments",
        "content_status",
        "created",
        "identifier",
        "keywords",
        "language",
        "last_modified_by",
        "last_printed",
        "modified",
        "revision",
        "subject",
        "title",
        "version",
    )

    # ---- Detection / meta ----------------------------------------------------

    def detect(self, path: Path) -> bool:
        return path.suffix.lower() in self.extensions

    def read_meta(self, path: Path) -> dict[str, Any]:
        stat = path.stat()
        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise HandlerError(f"failed to open pptx: {exc}") from exc

        slide_titles: list[str] = []
        slide_count = 0
        for slide in prs.slides:
            slide_count += 1
            slide_titles.append(_slide_title(slide))

        return {
            "format": self.name,
            "size": stat.st_size,
            "mtime": stat.st_mtime,
            "slide_count": slide_count,
            "slide_titles": slide_titles,
        }

    # ---- Text ----------------------------------------------------------------

    def read_text(self, path: Path, **_: Any) -> str:
        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise HandlerError(f"failed to open pptx: {exc}") from exc

        parts: list[str] = []
        for i, slide in enumerate(prs.slides):
            joined = "\n".join(_iter_text_frames(slide))
            parts.append(f"# Slide {i}\n{joined}\n\n")
        return "".join(parts)

    def write_text(self, path: Path, content: str, **_: Any) -> None:
        raise UnsupportedFormatError("write_text not supported for pptx")

    def extract_for_search(self, path: Path) -> str:
        try:
            prs = Presentation(str(path))
        except Exception:
            return ""
        chunks: list[str] = []
        for slide in prs.slides:
            chunks.extend(_iter_text_frames(slide))
        return "\n".join(chunks)

    # ---- Structured ops ------------------------------------------------------

    def structured_get(self, path: Path, expr: str) -> Any:
        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise HandlerError(f"failed to open pptx: {exc}") from exc

        if expr == "slides":
            out: list[dict[str, Any]] = []
            for i, slide in enumerate(prs.slides):
                texts = _iter_text_frames(slide)
                out.append(
                    {
                        "index": i,
                        "title": _slide_title(slide),
                        "text_count": len(texts),
                    }
                )
            return out

        if expr.startswith("slide:"):
            idx = _parse_slide_index(expr.split(":", 1)[1])
            slide = _get_slide(prs, idx)
            shapes = [_shape_to_dict(s) for s in slide.shapes]
            return {
                "index": idx,
                "title": _slide_title(slide),
                "shapes": shapes,
            }

        if expr.startswith("slide_text:"):
            idx = _parse_slide_index(expr.split(":", 1)[1])
            slide = _get_slide(prs, idx)
            return "\n".join(_iter_text_frames(slide))

        raise HandlerError(f"unsupported structured_get expression: {expr!r}")

    def structured_set(self, path: Path, expr: str, value: Any) -> None:
        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise HandlerError(f"failed to open pptx: {exc}") from exc

        if expr.startswith("slide_text:"):
            # Form: slide_text:<n>:<placeholder_name_or_idx>
            rest = expr.split(":", 1)[1]
            parts = rest.split(":", 1)
            if len(parts) != 2:
                raise HandlerError(
                    "slide_text set requires 'slide_text:<n>:<placeholder_name_or_idx>'"
                )
            idx = _parse_slide_index(parts[0])
            placeholder_key = parts[1]
            slide = _get_slide(prs, idx)

            if not isinstance(value, str):
                raise HandlerError("value must be a string for slide_text set")

            target = None
            # Try placeholder.idx match first (numeric).
            try:
                key_int = int(placeholder_key)
            except (TypeError, ValueError):
                key_int = None

            for ph in slide.placeholders:
                try:
                    ph_idx = ph.placeholder_format.idx
                except Exception:
                    ph_idx = None
                if key_int is not None and ph_idx == key_int:
                    target = ph
                    break
                if ph.name == placeholder_key:
                    target = ph
                    break

            if target is None:
                raise HandlerError(f"placeholder not found: {placeholder_key!r} on slide {idx}")
            try:
                if not target.has_text_frame:
                    raise HandlerError(f"placeholder {placeholder_key!r} has no text frame")
                target.text_frame.text = value
            except HandlerError:
                raise
            except Exception as exc:
                raise HandlerError(f"failed to set placeholder text: {exc}") from exc

            prs.save(str(path))
            return

        if expr.startswith("slide_title:"):
            idx = _parse_slide_index(expr.split(":", 1)[1])
            slide = _get_slide(prs, idx)
            if not isinstance(value, str):
                raise HandlerError("value must be a string for slide_title set")
            try:
                title_shape = slide.shapes.title
            except Exception as exc:
                raise HandlerError(f"slide has no title placeholder: {exc}") from exc
            if title_shape is None:
                raise HandlerError(f"slide {idx} has no title placeholder")
            try:
                title_shape.text_frame.text = value
            except Exception as exc:
                raise HandlerError(f"failed to set title: {exc}") from exc
            prs.save(str(path))
            return

        raise HandlerError(f"unsupported structured_set expression: {expr!r}")

    def structured_delete(self, path: Path, expr: str) -> None:
        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise HandlerError(f"failed to open pptx: {exc}") from exc

        if expr.startswith("slide:"):
            idx = _parse_slide_index(expr.split(":", 1)[1])
            xml_slides = prs.slides._sldIdLst  # type: ignore[attr-defined]
            slides_list = list(xml_slides)
            if idx < 0 or idx >= len(slides_list):
                raise HandlerError(
                    f"slide index out of range: {idx} (slide_count={len(slides_list)})"
                )
            xml_slides.remove(slides_list[idx])
            prs.save(str(path))
            return

        raise HandlerError(f"unsupported structured_delete expression: {expr!r}")

    # ---- metadata write ------------------------------------------------------

    def write_meta(self, path: Path, updates: dict[str, Any]) -> dict[str, Any]:
        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise HandlerError(f"failed to open pptx: {exc}") from exc
        cp = prs.core_properties
        applied: dict[str, Any] = {}
        for key, value in updates.items():
            if key not in self._WRITABLE_CORE_PROPS:
                raise HandlerError(
                    f"unknown core property: {key!r} "
                    f"(allowed: {list(self._WRITABLE_CORE_PROPS)})"
                )
            try:
                setattr(cp, key, value)
            except (AttributeError, TypeError, ValueError) as exc:
                raise HandlerError(
                    f"failed to set core property {key!r}: {exc}"
                ) from exc
            applied[key] = value
        try:
            prs.save(str(path))
        except Exception as exc:  # noqa: BLE001
            raise HandlerError(f"failed to save pptx: {exc}") from exc
        return applied

    def strip_meta(self, path: Path) -> dict[str, Any]:
        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise HandlerError(f"failed to open pptx: {exc}") from exc
        cp = prs.core_properties
        cleared: list[str] = []
        for key in self._WRITABLE_CORE_PROPS:
            try:
                current = getattr(cp, key, None)
                if isinstance(current, str):
                    setattr(cp, key, "")
                else:
                    setattr(cp, key, None)
                cleared.append(key)
            except (AttributeError, TypeError, ValueError):  # pragma: no cover
                # Defensive: matches the docx strip_meta pattern; not
                # reachable on the python-pptx version exercised by tests.
                continue
        try:
            prs.save(str(path))
        except Exception as exc:  # noqa: BLE001
            raise HandlerError(f"failed to save pptx: {exc}") from exc
        return {"stripped": cleared}


# Runtime-checkable protocol sanity assertion.
_handler: FormatHandler = PptxHandler()
default_registry.register(_handler)
