"""Structured (format-aware) batch find/replace for binary container formats.

``batch_replace_content`` deliberately refuses DOCX/XLSX/PPTX because raw
byte-level find/replace would corrupt the ZIP container. This module
provides a *safe* alternative that walks the document's structured content
(paragraphs / table cells / spreadsheet cells / slide text frames) and
writes the modified document back through the format's native writer.

Currently supported: ``.docx``, ``.xlsx``, ``.xlsm``, ``.pptx``.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from mcp.server.fastmcp import FastMCP

from ..context import AppContext
from ..errors import UnsupportedFormatError, ValidationError
from ..utils.locks import file_lock
from .batch import _iter_writable_files  # reuse glob iteration


_SUPPORTED_FORMATS = frozenset({"docx", "xlsx", "pptx"})

# Scope shape per format. Validated up-front so callers get a clear error
# when they pass a key that the underlying replacer can't act on.
_DOCX_SCOPE_KEYS = frozenset(
    {
        "headings_only",
        "tables_only",
        "paragraph_range",
        "heading_section",
        "exclude_styles",
        "include_styles",
    }
)
_XLSX_SCOPE_KEYS = frozenset({"sheets", "cell_range"})
_PPTX_SCOPE_KEYS = frozenset({"slides"})


def _validate_scope(scope: dict[str, Any] | None, fmt: str) -> dict[str, Any]:
    """Return a normalised scope dict. Empty/None = match everything."""
    if scope is None:
        return {}
    if not isinstance(scope, dict):
        raise ValidationError("scope must be a dict")
    allowed = {
        "docx": _DOCX_SCOPE_KEYS,
        "xlsx": _XLSX_SCOPE_KEYS,
        "pptx": _PPTX_SCOPE_KEYS,
    }.get(fmt, frozenset())
    unknown = set(scope) - allowed
    if unknown:
        raise ValidationError(
            f"unsupported scope keys for {fmt}: {sorted(unknown)} (allowed: {sorted(allowed)})"
        )
    return scope


def _docx_paragraph_in_section(doc: Any, idx: int, section: str) -> bool:
    """True if paragraph ``idx`` falls under a heading whose text matches ``section``."""
    rx = re.compile(section)
    current_level: int | None = None
    in_section = False
    for i, para in enumerate(doc.paragraphs):
        style_name = (getattr(getattr(para, "style", None), "name", None) or "").strip()
        m = re.match(r"^[Hh]eading\s+(\d+)$", style_name)
        if m:
            level = int(m.group(1))
            text = (para.text or "").strip()
            if rx.search(text):
                current_level = level
                in_section = True
            elif current_level is not None and level <= current_level:  # pragma: no branch
                in_section = False
                current_level = None
        if i == idx:
            return in_section
    return False  # pragma: no cover - idx must always exist within the loop


def _docx_paragraph_passes_scope(doc: Any, idx: int, para: Any, scope: dict[str, Any]) -> bool:
    """Return True if the paragraph at ``idx`` should be touched."""
    if not scope:
        return True
    if scope.get("tables_only"):
        return False  # paragraphs are skipped when only tables are in scope
    style_name = (getattr(getattr(para, "style", None), "name", None) or "").strip()
    if scope.get("headings_only"):
        if not re.match(r"^[Hh]eading\s+\d+$", style_name):
            return False
    rng = scope.get("paragraph_range")
    if rng is not None:
        if not (int(rng[0]) <= idx <= int(rng[1])):
            return False
    inc = scope.get("include_styles")
    if inc is not None and style_name not in set(inc):
        return False
    exc = scope.get("exclude_styles")
    if exc is not None and style_name in set(exc):
        return False
    section = scope.get("heading_section")
    if section is not None and not _docx_paragraph_in_section(doc, idx, section):
        return False
    return True


def _docx_table_cell_passes_scope(scope: dict[str, Any]) -> bool:
    """Tables are always in scope unless `headings_only` excludes them."""
    if not scope:
        return True
    if scope.get("headings_only"):
        return False
    return True


def _replace_in_docx(
    path: Path,
    pattern: re.Pattern[str],
    repl: str,
    *,
    apply: bool,
    scope: dict[str, Any] | None = None,
) -> tuple[int, list[dict[str, Any]]]:
    from docx import Document

    doc = Document(str(path))
    scope = scope or {}
    total = 0
    locations: list[dict[str, Any]] = []

    for idx, para in enumerate(doc.paragraphs):
        if not para.text:
            continue
        if not _docx_paragraph_passes_scope(doc, idx, para, scope):
            continue
        new_text, n = pattern.subn(repl, para.text)
        if n > 0:
            total += n
            locations.append({"kind": "paragraph", "index": idx, "matches": n})
            if apply:
                # Replace the entire paragraph text - preserves the
                # paragraph's style; per-run formatting inside is lost
                # (rare for plain text replace, acceptable trade-off).
                for run in list(para.runs):
                    run.text = ""
                if para.runs:
                    para.runs[0].text = new_text
                else:  # pragma: no cover — defensive: a python-docx
                    # paragraph that has text without runs is constructible
                    # only via raw OXML, never via the normal authoring path.
                    para.add_run(new_text)

    if _docx_table_cell_passes_scope(scope):
        for t_idx, table in enumerate(doc.tables):
            for r_idx, row in enumerate(table.rows):
                for c_idx, cell in enumerate(row.cells):
                    if not cell.text:
                        continue
                    new_text, n = pattern.subn(repl, cell.text)
                    if n > 0:
                        total += n
                        locations.append(
                            {
                                "kind": "table_cell",
                                "table": t_idx,
                                "row": r_idx,
                                "col": c_idx,
                                "matches": n,
                            }
                        )
                        if apply:
                            cell.text = new_text

    if apply and total > 0:
        doc.save(str(path))

    return total, locations


def _xlsx_cell_in_range(coord: str, ref_range: str) -> bool:
    """True if ``coord`` (e.g. ``B2``) falls inside ``ref_range`` (e.g. ``A1:E100``)."""
    from openpyxl.utils import column_index_from_string
    from openpyxl.utils.cell import coordinate_from_string

    parts = ref_range.split(":")
    if len(parts) != 2:
        raise ValidationError(f"invalid cell_range: {ref_range!r}")
    try:
        c0_letter, c0_row = coordinate_from_string(parts[0])
        c1_letter, c1_row = coordinate_from_string(parts[1])
        cur_letter, cur_row = coordinate_from_string(coord)
    except Exception as exc:  # noqa: BLE001 - openpyxl uses CellCoordinatesException
        raise ValidationError(f"invalid cell coordinate: {exc}") from exc
    c0 = column_index_from_string(c0_letter)
    c1 = column_index_from_string(c1_letter)
    cur = column_index_from_string(cur_letter)
    return min(c0, c1) <= cur <= max(c0, c1) and min(c0_row, c1_row) <= cur_row <= max(
        c0_row, c1_row
    )


def _replace_in_xlsx(
    path: Path,
    pattern: re.Pattern[str],
    repl: str,
    *,
    apply: bool,
    scope: dict[str, Any] | None = None,
) -> tuple[int, list[dict[str, Any]]]:
    import openpyxl

    wb = openpyxl.load_workbook(str(path))
    scope = scope or {}
    total = 0
    locations: list[dict[str, Any]] = []
    sheet_filter = set(scope["sheets"]) if scope.get("sheets") is not None else None
    cell_range = scope.get("cell_range")
    try:
        for ws in wb.worksheets:
            if sheet_filter is not None and ws.title not in sheet_filter:
                continue
            for row in ws.iter_rows():
                for cell in row:
                    val = cell.value
                    if not isinstance(val, str) or not val:
                        continue
                    if cell_range is not None and not _xlsx_cell_in_range(
                        cell.coordinate, cell_range
                    ):
                        continue
                    new_val, n = pattern.subn(repl, val)
                    if n > 0:
                        total += n
                        locations.append(
                            {
                                "kind": "cell",
                                "sheet": ws.title,
                                "ref": cell.coordinate,
                                "matches": n,
                            }
                        )
                        if apply:
                            cell.value = new_val
        if apply and total > 0:
            wb.save(str(path))
    finally:
        wb.close()

    return total, locations


def _replace_in_pptx(
    path: Path,
    pattern: re.Pattern[str],
    repl: str,
    *,
    apply: bool,
    scope: dict[str, Any] | None = None,
) -> tuple[int, list[dict[str, Any]]]:
    from pptx import Presentation

    prs = Presentation(str(path))
    scope = scope or {}
    total = 0
    locations: list[dict[str, Any]] = []
    slide_filter = set(int(i) for i in scope["slides"]) if scope.get("slides") is not None else None

    for s_idx, slide in enumerate(prs.slides):
        if slide_filter is not None and s_idx not in slide_filter:
            continue
        for sh_idx, shape in enumerate(slide.shapes):
            tf = getattr(shape, "text_frame", None)
            if tf is None:
                continue
            for p_idx, para in enumerate(tf.paragraphs):
                full_text = "".join(run.text for run in para.runs)
                if not full_text:
                    continue
                new_text, n = pattern.subn(repl, full_text)
                if n > 0:
                    total += n
                    locations.append(
                        {
                            "kind": "slide_text",
                            "slide": s_idx,
                            "shape": sh_idx,
                            "paragraph": p_idx,
                            "matches": n,
                        }
                    )
                    if apply:
                        runs = list(para.runs)
                        for run in runs[1:]:
                            run.text = ""
                        if runs:
                            runs[0].text = new_text
                        else:  # pragma: no cover — defensive: a paragraph
                            # with non-empty .text but zero runs is not
                            # producible via the python-pptx public API.
                            new_run = para.add_run()
                            new_run.text = new_text

    if apply and total > 0:
        prs.save(str(path))

    return total, locations


def register(mcp: FastMCP, ctx: AppContext) -> None:
    @mcp.tool(
        name="batch_replace_structured",
        description=(
            "Format-aware find/replace inside DOCX/XLSX/PPTX files. Walks "
            "paragraphs, table cells, spreadsheet cells, and slide text "
            "frames via the native writer (no raw bytes), so the binary "
            "container stays valid. `regex=False` and `case_sensitive=True` "
            "by default. Always snapshots pre+post when applying. "
            "Pass `scope` to restrict the replace to a subset of the doc: "
            "DOCX accepts {headings_only, tables_only, paragraph_range, "
            "heading_section, exclude_styles, include_styles}; "
            "XLSX accepts {sheets, cell_range}; "
            "PPTX accepts {slides}."
        ),
    )
    def batch_replace_structured(
        glob: str,
        old: str,
        new: str,
        regex: bool = False,
        dry_run: bool = True,
        case_sensitive: bool = True,
        scope: dict[str, Any] | None = None,
    ) -> dict[str, Any]:
        flags = 0 if case_sensitive else re.IGNORECASE
        pattern = re.compile(old if regex else re.escape(old), flags)

        plan: list[dict[str, Any]] = []
        skipped: list[dict[str, Any]] = []

        for root_name, p, root_abs in _iter_writable_files(ctx, glob):
            uri = f"{root_name}:/{p.relative_to(root_abs).as_posix()}"
            handler = ctx.registry.for_path(p)
            if handler is None or handler.name not in _SUPPORTED_FORMATS:
                skipped.append({"uri": uri, "reason": "format_not_supported"})
                continue

            # Validate scope per-format up-front for a clear error message.
            try:
                normalised_scope = _validate_scope(scope, handler.name)
            except ValidationError:
                raise

            # Step 1 — always do a no-write pass first so we can pre-snapshot
            # before mutating, and so dry runs are completely side-effect-free.
            try:
                if handler.name == "docx":
                    total, locs = _replace_in_docx(
                        p, pattern, new, apply=False, scope=normalised_scope
                    )
                elif handler.name == "xlsx":
                    total, locs = _replace_in_xlsx(
                        p, pattern, new, apply=False, scope=normalised_scope
                    )
                else:  # pptx
                    total, locs = _replace_in_pptx(
                        p, pattern, new, apply=False, scope=normalised_scope
                    )
            except UnsupportedFormatError:
                raise
            except ValidationError:
                raise
            except Exception as exc:  # noqa: BLE001
                skipped.append({"uri": uri, "reason": "render_failed", "error": str(exc)})
                continue

            if total == 0:
                continue

            entry: dict[str, Any] = {
                "uri": uri,
                "absolute": str(p),
                "format": handler.name,
                "replacements": total,
                "locations": locs,
            }
            plan.append(entry)

            if dry_run:
                continue

            rel = p.relative_to(root_abs).as_posix()
            with file_lock(p):
                # Pre-snapshot before mutation so we can roll back.
                try:
                    ctx.versions.snapshot(
                        root_name=root_name,
                        rel_path=rel,
                        source=p,
                        action="batch_replace_structured_pre",
                    )
                except Exception:  # noqa: BLE001
                    pass
                # Step 2 — apply.
                try:
                    if handler.name == "docx":
                        _replace_in_docx(p, pattern, new, apply=True, scope=normalised_scope)
                    elif handler.name == "xlsx":
                        _replace_in_xlsx(p, pattern, new, apply=True, scope=normalised_scope)
                    else:  # pptx
                        _replace_in_pptx(p, pattern, new, apply=True, scope=normalised_scope)
                except Exception as exc:  # noqa: BLE001
                    # Demote planned entry into skipped on apply failure.
                    plan.pop()
                    skipped.append({"uri": uri, "reason": "apply_failed", "error": str(exc)})
                    continue
                try:
                    ctx.versions.snapshot(
                        root_name=root_name,
                        rel_path=rel,
                        source=p,
                        action="batch_replace_structured_post",
                    )
                except Exception:  # noqa: BLE001
                    pass
            ctx.audit.log(
                "batch_replace_structured",
                uri=uri,
                format=handler.name,
                replacements=total,
            )

        result: dict[str, Any] = {
            "dry_run": dry_run,
            "count": len(plan),
            "files": plan,
        }
        if skipped:
            result["skipped"] = skipped
            summary: dict[str, int] = {}
            for s in skipped:
                summary[s["reason"]] = summary.get(s["reason"], 0) + 1
            result["skipped_summary"] = summary
        return result
