"""Tests for :mod:`dokumen_pintar.tools.batch_structured`."""

from __future__ import annotations

from pathlib import Path
from typing import Callable

import pytest
from mcp.server.fastmcp import FastMCP

from dokumen_pintar.config import AppConfig
from dokumen_pintar.context import build_context
from dokumen_pintar.tools import batch_structured


def _setup(cfg: AppConfig):  # type: ignore[no-untyped-def]
    ctx = build_context(cfg)
    mcp = FastMCP(name="test")
    batch_structured.register(mcp, ctx)
    return mcp, ctx


def _tool(mcp: FastMCP, name: str):  # type: ignore[no-untyped-def]
    return mcp._tool_manager._tools[name].fn


def _make_docx(path: Path, paragraphs: list[str]) -> None:
    from docx import Document

    doc = Document()
    for para in paragraphs:
        doc.add_paragraph(para)
    doc.save(str(path))


def _make_docx_with_table(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_paragraph("Hello world")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "world A"
    table.rows[0].cells[1].text = "B"
    table.rows[1].cells[0].text = "world C"
    table.rows[1].cells[1].text = "D"
    doc.save(str(path))


def _make_xlsx(path: Path, values: list[list[str]]) -> None:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    for r_idx, row in enumerate(values, start=1):
        for c_idx, val in enumerate(row, start=1):
            ws.cell(row=r_idx, column=c_idx, value=val)
    wb.save(str(path))


def _make_pptx(path: Path, texts: list[str]) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[5]
    for text in texts:
        slide = prs.slides.add_slide(blank)
        tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tx.text_frame.text = text
    prs.save(str(path))


def test_dry_run_docx_paragraph(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    target = docs_dir / "doc.docx"
    _make_docx(target, ["Hello world", "no match here", "another world"])
    original = target.read_bytes()
    mcp, _ = _setup(make_config())
    result = _tool(mcp, "batch_replace_structured")(
        glob="*.docx", old="world", new="WORLD", dry_run=True
    )
    assert result["dry_run"] is True
    assert result["count"] == 1
    file_entry = result["files"][0]
    assert file_entry["replacements"] == 2
    assert file_entry["format"] == "docx"
    # Dry run must NOT touch the file at all.
    assert target.read_bytes() == original


def test_apply_docx_paragraph_and_table(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    from docx import Document

    docs_dir, _ = tmp_roots
    target = docs_dir / "doc.docx"
    _make_docx_with_table(target)
    mcp, _ = _setup(make_config())
    result = _tool(mcp, "batch_replace_structured")(
        glob="*.docx", old="world", new="WORLD", dry_run=False
    )
    assert result["count"] == 1
    # Re-open and verify text replaced.
    doc = Document(str(target))
    text = "\n".join(p.text for p in doc.paragraphs)
    assert "WORLD" in text
    cell_text = doc.tables[0].rows[0].cells[0].text
    assert cell_text == "WORLD A"


def test_dry_run_xlsx(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    target = docs_dir / "data.xlsx"
    _make_xlsx(target, [["foo", "bar"], ["foobar", "qux"]])
    original = target.read_bytes()
    mcp, _ = _setup(make_config())
    result = _tool(mcp, "batch_replace_structured")(
        glob="*.xlsx", old="foo", new="FOO", dry_run=True
    )
    assert result["count"] == 1
    assert result["files"][0]["replacements"] == 2
    assert target.read_bytes() == original


def test_apply_xlsx(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    import openpyxl

    docs_dir, _ = tmp_roots
    target = docs_dir / "data.xlsx"
    _make_xlsx(target, [["foo here", "bar"], ["another foo", "x"]])
    mcp, _ = _setup(make_config())
    _tool(mcp, "batch_replace_structured")(
        glob="*.xlsx", old="foo", new="FOO", dry_run=False
    )
    wb = openpyxl.load_workbook(str(target))
    try:
        ws = wb.active
        assert ws.cell(row=1, column=1).value == "FOO here"
        assert ws.cell(row=2, column=1).value == "another FOO"
    finally:
        wb.close()


def test_apply_pptx(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    from pptx import Presentation

    docs_dir, _ = tmp_roots
    target = docs_dir / "deck.pptx"
    _make_pptx(target, ["Hello world", "world peace"])
    mcp, _ = _setup(make_config())
    result = _tool(mcp, "batch_replace_structured")(
        glob="*.pptx", old="world", new="WORLD", dry_run=False
    )
    assert result["count"] == 1
    prs = Presentation(str(target))
    seen = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                seen.append(shape.text_frame.text)
    assert any("WORLD" in t for t in seen)


def test_unsupported_format_skipped(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    (docs_dir / "plain.txt").write_text("hello", encoding="utf-8")
    mcp, _ = _setup(make_config())
    result = _tool(mcp, "batch_replace_structured")(
        glob="*.txt", old="hello", new="hi", dry_run=True
    )
    assert result["count"] == 0
    assert result.get("skipped_summary", {}).get("format_not_supported") == 1


def test_no_match_returns_empty_plan(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    target = docs_dir / "doc.docx"
    _make_docx(target, ["Hello"])
    mcp, _ = _setup(make_config())
    result = _tool(mcp, "batch_replace_structured")(
        glob="*.docx", old="ZZZNOPE", new="x", dry_run=True
    )
    assert result["count"] == 0


def test_regex_mode(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    from docx import Document

    docs_dir, _ = tmp_roots
    target = docs_dir / "doc.docx"
    _make_docx(target, ["item123 and item456"])
    mcp, _ = _setup(make_config())
    _tool(mcp, "batch_replace_structured")(
        glob="*.docx", old=r"item\d+", new="ITEM", regex=True, dry_run=False
    )
    doc = Document(str(target))
    text = "\n".join(p.text for p in doc.paragraphs)
    assert text.count("ITEM") == 2


def test_case_insensitive(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    from docx import Document

    docs_dir, _ = tmp_roots
    target = docs_dir / "doc.docx"
    _make_docx(target, ["Hello HELLO hello"])
    mcp, _ = _setup(make_config())
    _tool(mcp, "batch_replace_structured")(
        glob="*.docx",
        old="hello",
        new="HI",
        case_sensitive=False,
        dry_run=False,
    )
    doc = Document(str(target))
    text = "\n".join(p.text for p in doc.paragraphs)
    assert text == "HI HI HI"


def test_corrupt_docx_recorded_as_render_failed(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    target = docs_dir / "broken.docx"
    target.write_bytes(b"not a real docx")
    mcp, _ = _setup(make_config())
    result = _tool(mcp, "batch_replace_structured")(
        glob="*.docx", old="x", new="y", dry_run=False
    )
    assert result["count"] == 0
    reasons = {s["reason"] for s in result.get("skipped", [])}
    assert "render_failed" in reasons


def test_apply_failure_demoted_to_skipped(
    make_config: Callable[..., AppConfig],
    tmp_roots: tuple[Path, Path],
    monkeypatch,
) -> None:
    """If the apply pass crashes after the dry pass succeeded, the file is
    moved from ``files`` into ``skipped`` with reason ``apply_failed``."""
    docs_dir, _ = tmp_roots
    target = docs_dir / "doc.docx"
    _make_docx(target, ["Hello world"])
    mcp, _ = _setup(make_config())

    real_replace = batch_structured._replace_in_docx
    state = {"calls": 0}

    def flaky(path, pattern, repl, *, apply: bool):  # type: ignore[no-untyped-def]
        state["calls"] += 1
        if apply:
            raise RuntimeError("disk full")
        return real_replace(path, pattern, repl, apply=apply)

    monkeypatch.setattr(batch_structured, "_replace_in_docx", flaky)

    result = _tool(mcp, "batch_replace_structured")(
        glob="*.docx", old="world", new="X", dry_run=False
    )
    assert result["count"] == 0
    reasons = {s["reason"] for s in result.get("skipped", [])}
    assert "apply_failed" in reasons


def test_glob_root_prefix_works(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    _make_docx(docs_dir / "a.docx", ["world"])
    mcp, _ = _setup(make_config())
    result = _tool(mcp, "batch_replace_structured")(
        glob="documents:/*.docx", old="world", new="X", dry_run=True
    )
    assert result["count"] == 1
