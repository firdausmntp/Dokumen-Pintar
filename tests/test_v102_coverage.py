"""v1.0.2 coverage-completion tests.

Targets specific defensive branches and edge paths that the format-level
test suites do not exercise on their own. Each test in this file
corresponds to one or more concrete "Missing" lines reported by
``pytest --cov-report=term-missing``.

The file is intentionally consolidated rather than split per-module so we
can keep the v1.0.2 coverage gate (``fail_under = 100``) green with a
single, scannable artifact.
"""

from __future__ import annotations

from pathlib import Path
from typing import Callable, Iterator
from unittest.mock import patch

import pytest
from mcp.server.fastmcp import FastMCP

from dokumen_pintar.authoring import SpecError, validate_spec
from dokumen_pintar.authoring import markdown_to_spec as md2spec_mod
from dokumen_pintar.authoring import render_docx as render_docx_mod
from dokumen_pintar.authoring import render_pdf as render_pdf_mod
from dokumen_pintar.authoring.markdown_to_spec import markdown_to_spec
from dokumen_pintar.authoring.render_docx import render_docx
from dokumen_pintar.authoring.render_pdf import render_pdf
from dokumen_pintar.authoring.markdown_to_spec import (
    _inline_to_runs,
    _list_from_tokens,
    _table_from_tokens,
)
from dokumen_pintar.authoring.spec import _validate_run, _validate_image, validate_spec as _validate_spec
from dokumen_pintar.config import AppConfig
from dokumen_pintar.context import build_context
from dokumen_pintar.errors import HandlerError, UnsupportedFormatError, ValidationError
from dokumen_pintar.handlers import latex_handler as latex_mod
from dokumen_pintar.handlers.latex_handler import (
    LatexHandler,
    _collect_environments,
    _collect_outline,
    _macro_argument_text,
)
from dokumen_pintar.tools import authoring as authoring_tool
from dokumen_pintar.tools import batch_structured
from dokumen_pintar.tools import search as search_tool
from dokumen_pintar.tools import structured, version, workspace, content_crud, file_crud


# ─────────────────────────── helpers


def _setup(cfg: AppConfig, *, modules):  # type: ignore[no-untyped-def]
    ctx = build_context(cfg)
    mcp = FastMCP(name="cov-test")
    workspace.register(mcp, ctx)
    file_crud.register(mcp, ctx)
    content_crud.register(mcp, ctx)
    structured.register(mcp, ctx)
    version.register(mcp, ctx)
    for mod in modules:
        mod.register(mcp, ctx)
    return mcp, ctx


def _tool(mcp, name):  # type: ignore[no-untyped-def]
    return mcp._tool_manager._tools[name].fn


# ─────────────────────────── search.py — conflicting glob root


class TestSearchConflictingRoot:
    def test_conflicting_glob_and_filter_yields_nothing(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        docs_dir, ref_dir = tmp_roots
        (docs_dir / "a.txt").write_text("hello")
        (ref_dir / "b.txt").write_text("hello")
        mcp, _ = _setup(make_config(), modules=[search_tool])
        # glob says "documents:/...", but the user passes root="ref" too.
        out = _tool(mcp, "search_filename")(
            glob_pattern="documents:/*.txt", root="ref"
        )
        # Conflict between URI prefix and explicit root → empty result.
        assert out["matches"] == []
        assert out["count"] == 0


# ─────────────────────────── tools/authoring.py


class TestAuthoringEdgeCases:
    def test_path_resolver_for_resolves_via_guard(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        """When the resolver succeeds via PathGuard we hit the absolute-path
        return — covers ``return r.absolute`` (line 45)."""
        docs_dir, _ = tmp_roots
        target = docs_dir / "real.png"
        target.write_bytes(b"\x89PNG\r\n\x1a\n")
        ctx = build_context(make_config())
        resolver = authoring_tool._path_resolver_for(ctx)
        out = resolver("documents:/real.png")
        assert out == target

    def test_snapshot_pre_swallows_exception(
        self, make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
    ) -> None:
        docs_dir, _ = tmp_roots
        target = docs_dir / "x.docx"
        target.write_bytes(b"placeholder")  # exists() True so snapshot runs
        ctx = build_context(make_config())
        with patch.object(
            ctx.versions, "snapshot", side_effect=RuntimeError("io")
        ):
            authoring_tool._snapshot_pre(
                ctx, target, "documents", "x.docx", "compose_pre"
            )

    def test_snapshot_post_swallows_exception_returns_none(
        self, make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
    ) -> None:
        docs_dir, _ = tmp_roots
        target = docs_dir / "x.docx"
        target.write_bytes(b"placeholder")
        ctx = build_context(make_config())
        with patch.object(
            ctx.versions, "snapshot", side_effect=RuntimeError("io")
        ):
            assert (
                authoring_tool._snapshot_post(
                    ctx, target, "documents", "x.docx", "compose_post"
                )
                is None
            )

    def test_compose_from_markdown_invalid_md_raises_validation_error(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        """If ``markdown_to_spec`` raises :class:`SpecError` the tool must
        re-raise as :class:`ValidationError` (lines 207-208)."""
        mcp, _ = _setup(make_config(), modules=[authoring_tool])
        with patch.object(
            authoring_tool,
            "markdown_to_spec",
            side_effect=SpecError("boom"),
        ):
            with pytest.raises(ValidationError):
                _tool(mcp, "compose_from_markdown")(
                    path="documents:/out.docx", markdown="# hi"
                )


# ─────────────────────────── authoring/spec.py validators


class TestSpecValidators:
    def test_run_must_be_dict(self) -> None:
        with pytest.raises(SpecError):
            _validate_run("not a dict", idx=0, block_idx=0)  # type: ignore[arg-type]

    def test_run_color_must_be_string(self) -> None:
        with pytest.raises(SpecError):
            _validate_run(
                {"text": "x", "color": 123}, idx=0, block_idx=0
            )

    def test_paragraph_runs_must_be_list(self) -> None:
        with pytest.raises(SpecError):
            _validate_spec({"blocks": [{"type": "paragraph", "runs": "nope"}]})

    def test_paragraph_text_must_be_string(self) -> None:
        with pytest.raises(SpecError):
            _validate_spec({"blocks": [{"type": "paragraph", "text": 42}]})

    def test_table_row_must_be_list(self) -> None:
        with pytest.raises(SpecError):
            _validate_spec(
                {"blocks": [{"type": "table", "rows": ["not", "rows"]}]}
            )

    def test_image_caption_must_be_string(self) -> None:
        with pytest.raises(SpecError):
            _validate_image({"path": "x.png", "caption": 5}, 0)

    def test_code_language_must_be_string(self) -> None:
        with pytest.raises(SpecError):
            _validate_spec(
                {"blocks": [{"type": "code", "text": "x", "language": 1}]}
            )


# ─────────────────────────── authoring/render_docx — branch


class TestRenderDocxBranches:
    def test_table_no_header_with_rows(self, tmp_path: Path) -> None:
        spec = validate_spec(
            {"blocks": [{"type": "table", "rows": [["1", "2"], ["3", "4"]]}]}
        )
        out = tmp_path / "noh.docx"
        render_docx_mod.render_docx(spec, out)
        from docx import Document

        doc = Document(str(out))
        assert len(doc.tables[0].rows) == 2

    def test_image_without_width_cm(self, tmp_path: Path) -> None:
        """Cover the ``width_cm`` False branch in render_docx (line 106→108)."""
        # Create a tiny PNG fixture.
        png = tmp_path / "tiny.png"
        from PIL import Image as PILImage

        PILImage.new("RGB", (4, 4), color=(255, 0, 0)).save(png, "PNG")
        spec = validate_spec(
            {"blocks": [{"type": "image", "path": str(png)}]}
        )
        out = tmp_path / "imgnodim.docx"
        render_docx_mod.render_docx(spec, out)
        assert out.exists() and out.stat().st_size > 0


# ─────────────────────────── authoring/render_pdf — line + branch


class TestRenderPdfBranches:
    def test_underline_run(self, tmp_path: Path) -> None:
        """Cover the underline run branch (line 79)."""
        spec = validate_spec(
            {
                "blocks": [
                    {
                        "type": "paragraph",
                        "runs": [{"text": "u", "underline": True}],
                    }
                ]
            }
        )
        out = tmp_path / "u.pdf"
        render_pdf(spec, out)
        assert out.exists() and out.read_bytes()[:4] == b"%PDF"

    def test_image_without_width_cm(self, tmp_path: Path) -> None:
        png = tmp_path / "tiny.png"
        from PIL import Image as PILImage

        PILImage.new("RGB", (4, 4), color=(0, 255, 0)).save(png, "PNG")
        spec = validate_spec(
            {"blocks": [{"type": "image", "path": str(png)}]}
        )
        out = tmp_path / "imgnodim.pdf"
        render_pdf(spec, out)
        assert out.exists() and out.read_bytes()[:4] == b"%PDF"


# ─────────────────────────── markdown_to_spec edge cases


class TestMarkdownToSpecEdges:
    def test_image_with_empty_alt_dropped(self) -> None:
        """`![](pic.png)` produces no marker run."""

        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content
                self.children: list = []

        inline = T("inline")
        inline.children = [T("image", "")]
        assert _inline_to_runs(inline) == []

    def test_unknown_inline_token_no_content(self) -> None:
        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content
                self.children: list = []

        inline = T("inline")
        inline.children = [T("totally_unknown", "")]
        assert _inline_to_runs(inline) == []

    def test_table_orphan_tr_close(self) -> None:
        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content

        tokens = [T("table_open"), T("tr_close"), T("table_close")]
        block, _ = _table_from_tokens(tokens, 0)
        assert block["rows"] == []

    def test_table_tr_close_outside_head_or_body(self) -> None:
        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content

        tokens = [
            T("table_open"),
            T("tr_open"),
            T("inline", "x"),
            T("tr_close"),
            T("table_close"),
        ]
        block, _ = _table_from_tokens(tokens, 0)
        assert block["header"] is None and block["rows"] == []

    def test_list_empty_inline_in_item(self) -> None:
        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content

        tokens = [
            T("bullet_list_open"),
            T("list_item_open"),
            T("inline", ""),
            T("list_item_close"),
            T("bullet_list_close"),
        ]
        block, _ = _list_from_tokens(tokens, 0, ordered=False)
        assert block["items"] == []

    def test_blockquote_at_eof(self) -> None:
        class FakeTok:
            def __init__(self, type_: str = "blockquote_open") -> None:
                self.type = type_
                self.content = ""
                self.tag = ""
                self.info = ""
                self.map = None
                self.children = None

        class FakeMD:
            def parse(self, src):  # type: ignore[no-untyped-def]
                return [FakeTok()]

        with patch.object(md2spec_mod, "_build_md", lambda: FakeMD()):
            spec = markdown_to_spec("ignored")
        bqs = [b for b in spec.blocks if b["type"] == "blockquote"]
        assert bqs and bqs[0]["text"] == ""

    def test_empty_bullet_list_dropped(self) -> None:
        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content
                self.tag = ""
                self.info = ""
                self.map = None
                self.children = None

        class FakeMD:
            def parse(self, src):  # type: ignore[no-untyped-def]
                return [T("bullet_list_open"), T("bullet_list_close")]

        with patch.object(md2spec_mod, "_build_md", lambda: FakeMD()):
            spec = markdown_to_spec("ignored")
        assert all(b["type"] != "list" for b in spec.blocks)

    def test_empty_ordered_list_dropped(self) -> None:
        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content
                self.tag = ""
                self.info = ""
                self.map = None
                self.children = None

        class FakeMD:
            def parse(self, src):  # type: ignore[no-untyped-def]
                return [T("ordered_list_open"), T("ordered_list_close")]

        with patch.object(md2spec_mod, "_build_md", lambda: FakeMD()):
            spec = markdown_to_spec("ignored")
        assert all(b["type"] != "list" for b in spec.blocks)

    def test_unknown_top_level_token_skipped(self) -> None:
        class FakeTok:
            def __init__(self) -> None:
                self.type = "totally_unknown_top_level"
                self.content = ""
                self.tag = ""
                self.info = ""
                self.map = None
                self.children = None

        class FakeMD:
            def parse(self, src):  # type: ignore[no-untyped-def]
                return [FakeTok()]

        with patch.object(md2spec_mod, "_build_md", lambda: FakeMD()):
            spec = markdown_to_spec("ignored")
        # Unknown tokens skipped — no blocks produced.
        assert spec.blocks == []

    def test_paragraph_with_only_newline_strip_to_empty(self) -> None:
        spec = markdown_to_spec("\n\n")
        assert spec.blocks == []

    def test_inline_code_run(self) -> None:
        spec = markdown_to_spec("Use `python` here")
        ps = [b for b in spec.blocks if b["type"] == "paragraph"]
        assert ps and any(r.get("code") for r in ps[0]["runs"])

    def test_ordered_list_basic(self) -> None:
        spec = markdown_to_spec("1. one\n2. two\n")
        lists = [b for b in spec.blocks if b["type"] == "list"]
        assert lists and lists[0]["ordered"] is True

    def test_nested_blockquote(self) -> None:
        spec = markdown_to_spec("> outer\n>\n> > inner\n")
        bqs = [b for b in spec.blocks if b["type"] == "blockquote"]
        assert bqs

    def test_hr_block(self) -> None:
        spec = markdown_to_spec("text\n\n---\n\nmore\n")
        assert any(b["type"] == "hr" for b in spec.blocks)

    def test_fenced_code_block(self) -> None:
        spec = markdown_to_spec("```python\nprint(1)\n```\n")
        codes = [b for b in spec.blocks if b["type"] == "code"]
        assert codes and codes[0]["language"] == "python"

    def test_math_block(self) -> None:
        spec = markdown_to_spec("$$x = y$$\n")
        # Either captured as math or (depending on plugin presence) as
        # paragraph/code; just ensure no crash.
        assert spec.blocks

    def test_softbreak_emits_space_run(self) -> None:
        """Two consecutive lines joined by a softbreak (line 48)."""

        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content
                self.children: list = []

        inline = T("inline")
        inline.children = [T("text", "a"), T("softbreak", ""), T("text", "b")]
        runs = _inline_to_runs(inline)
        # The softbreak becomes a {"text": " "} run, then coalesced into the run list.
        assert any(r["text"] == " " for r in runs)

    def test_strikethrough_tokens_dropped(self) -> None:
        """Cover ``s_open`` / ``s_close`` pass branches (lines 58-61)."""

        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content
                self.children: list = []

        inline = T("inline")
        inline.children = [
            T("text", "a"),
            T("s_open", ""),
            T("text", "struck"),
            T("s_close", ""),
            T("text", "b"),
        ]
        runs = _inline_to_runs(inline)
        # Strike-through tokens are no-ops; surrounding text flows through.
        joined = "".join(r["text"] for r in runs)
        assert "struck" in joined

    def test_html_inline_dropped(self) -> None:
        """Cover ``html_inline`` pass branch (line 74)."""

        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content
                self.children: list = []

        inline = T("inline")
        inline.children = [
            T("text", "before"),
            T("html_inline", "<span>"),
            T("text", "after"),
        ]
        runs = _inline_to_runs(inline)
        joined = "".join(r["text"] for r in runs)
        assert "<span>" not in joined

    def test_unknown_inline_with_content_becomes_text(self) -> None:
        """Cover the fallback append path (line 79)."""

        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content
                self.children: list = []

        inline = T("inline")
        inline.children = [T("totally_unknown", "literal content")]
        runs = _inline_to_runs(inline)
        assert any(r["text"] == "literal content" for r in runs)

    def test_table_unterminated_returns_what_we_have(self) -> None:
        """Cover line 118 — fall-through return when table_close never appears."""

        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content

        tokens = [T("table_open"), T("tr_open"), T("tr_close")]  # no table_close
        block, _ = _table_from_tokens(tokens, 0)
        assert block["type"] == "table"

    def test_list_nested_increases_depth(self) -> None:
        """Cover line 137 — nested list_open inside item bumps depth."""

        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content

        tokens = [
            T("bullet_list_open"),
            T("list_item_open"),
            T("inline", "outer"),
            T("bullet_list_open"),  # nested
            T("list_item_open"),
            T("inline", "inner"),
            T("list_item_close"),
            T("bullet_list_close"),
            T("list_item_close"),
            T("bullet_list_close"),
        ]
        block, _ = _list_from_tokens(tokens, 0, ordered=False)
        # Outer item collected (inner is flattened as separate item).
        assert block["items"]

    def test_list_outer_close_with_depth_not_zero_continues(self) -> None:
        """Cover the 134→150 branch — outer ``bullet_list_close`` while
        depth != 0 must continue iteration."""

        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content

        # Open a nested list with an extra spurious open so depth > 1 when
        # the first close is encountered.
        tokens = [
            T("bullet_list_open"),
            T("bullet_list_open"),
            T("list_item_open"),
            T("inline", "x"),
            T("list_item_close"),
            T("bullet_list_close"),
            T("bullet_list_close"),
        ]
        block, _ = _list_from_tokens(tokens, 0, ordered=False)
        assert block["type"] == "list"

    def test_list_unterminated_returns_what_we_have(self) -> None:
        """Cover line 151 — fall-through return when close never appears."""

        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content

        tokens = [
            T("bullet_list_open"),
            T("list_item_open"),
            T("inline", "x"),
        ]
        block, _ = _list_from_tokens(tokens, 0, ordered=False)
        assert block["items"] == []  # never closed → never appended

    def test_paragraph_with_empty_inline_yields_empty_run(self) -> None:
        """A paragraph whose inline produces no runs falls back to a single
        empty-text run (line 182)."""

        class T:
            def __init__(self, type_: str, content: str = "") -> None:
                self.type = type_
                self.content = content
                self.tag = ""
                self.info = ""
                self.map = None
                self.children: list = []

        class FakeMD:
            def parse(self, src):  # type: ignore[no-untyped-def]
                inline = T("inline", "")
                inline.children = []
                return [T("paragraph_open"), inline, T("paragraph_close")]

        with patch.object(md2spec_mod, "_build_md", lambda: FakeMD()):
            spec = markdown_to_spec("ignored")
        paragraphs = [b for b in spec.blocks if b["type"] == "paragraph"]
        assert paragraphs and paragraphs[0]["runs"] == [{"text": ""}]


# ─────────────────────────── latex_handler edge cases


class TestLatexHandlerEdges:
    def test_macro_argument_no_args(self) -> None:
        """``_macro_argument_text`` returns "" when there are no args."""

        class FakeMacro:
            nodeargd = None

        assert _macro_argument_text(FakeMacro()) == ""

    def test_macro_argument_no_brace_group(self) -> None:
        """No mandatory ``{...}`` group → empty string."""

        class FakeArg:
            delimiters = ("[", "]")
            nodelist: list = []

        class _ArgD:
            argnlist = [None, FakeArg()]

        class FakeMacro:
            nodeargd = _ArgD()

        assert _macro_argument_text(FakeMacro()) == ""

    def test_collect_outline_parse_failure(self) -> None:
        from pylatexenc.latexwalker import LatexWalker

        with patch.object(
            LatexWalker, "get_latex_nodes",
            side_effect=RuntimeError("bad latex"),
        ):
            with pytest.raises(HandlerError):
                _collect_outline(r"\section{x}")

    def test_collect_environments_parse_failure(self) -> None:
        from pylatexenc.latexwalker import LatexWalker

        with patch.object(
            LatexWalker, "get_latex_nodes",
            side_effect=RuntimeError("bad latex"),
        ):
            with pytest.raises(HandlerError):
                _collect_environments(r"\begin{x}\end{x}")

    def test_read_meta_swallows_outline_failure(
        self, tmp_path: Path
    ) -> None:
        p = tmp_path / "x.tex"
        p.write_text(r"\documentclass{article}\section{Hi}")
        with patch.object(
            latex_mod, "_collect_outline",
            side_effect=HandlerError("nope"),
        ):
            meta = LatexHandler().read_meta(p)
        assert meta["outline"] == []

    def test_read_meta_swallows_envs_failure(self, tmp_path: Path) -> None:
        p = tmp_path / "x.tex"
        p.write_text(r"\documentclass{article}\begin{document}x\end{document}")
        with patch.object(
            latex_mod, "_collect_environments",
            side_effect=HandlerError("nope"),
        ):
            meta = LatexHandler().read_meta(p)
        assert meta["environment_counts"] == {}

    def test_extract_for_search_falls_back_to_raw_text(
        self, tmp_path: Path
    ) -> None:
        p = tmp_path / "x.tex"
        p.write_text(r"\section{Hello} world")
        from pylatexenc.latex2text import LatexNodes2Text

        with patch.object(
            LatexNodes2Text, "latex_to_text",
            side_effect=RuntimeError("boom"),
        ):
            out = LatexHandler().extract_for_search(p)
        # Fallback returns the raw text untouched.
        assert "Hello" in out


# ─────────────────────────── batch_structured edge cases


class TestBatchStructuredEdges:
    def _make_docx(self, path: Path, *, with_empty_para: bool = True) -> None:
        from docx import Document

        d = Document()
        if with_empty_para:
            d.add_paragraph("")  # explicitly empty paragraph
        d.add_paragraph("hit me")
        # Add a table with one empty cell + one non-empty cell.
        t = d.add_table(rows=1, cols=2)
        t.rows[0].cells[0].text = ""
        t.rows[0].cells[1].text = "hit me"
        d.save(str(path))

    def test_docx_empty_para_and_empty_cell_skipped(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        """Cover the ``if not para.text: continue`` (line 40) and
        ``if not cell.text: continue`` (line 62) branches in
        ``_replace_in_docx``."""
        docs_dir, _ = tmp_roots
        target = docs_dir / "mix.docx"
        self._make_docx(target)
        mcp, _ = _setup(make_config(), modules=[batch_structured])
        out = _tool(mcp, "batch_replace_structured")(
            glob="documents:/*.docx", old="hit me", new="X", dry_run=True
        )
        assert out["count"] == 1
        # 1 paragraph + 1 cell = 2 replacements (the empty ones skipped).
        assert out["files"][0]["replacements"] == 2

    def test_xlsx_non_string_cell_skipped(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        import openpyxl

        docs_dir, _ = tmp_roots
        target = docs_dir / "mix.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = 42
        ws["A2"] = ""
        ws["A3"] = "hit me"
        wb.save(str(target))
        mcp, _ = _setup(make_config(), modules=[batch_structured])
        out = _tool(mcp, "batch_replace_structured")(
            glob="documents:/*.xlsx", old="hit me", new="X", dry_run=True
        )
        assert out["count"] == 1

    def test_pptx_shape_without_text_frame(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        """A picture shape has no ``text_frame`` and must be skipped."""
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image as PILImage

        docs_dir, _ = tmp_roots
        png = docs_dir / "_pic.png"
        PILImage.new("RGB", (8, 8)).save(png, "PNG")
        target = docs_dir / "deck.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(str(png), Inches(1), Inches(1))
        slide.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1)).text_frame.text = (
            "hit me"
        )
        prs.save(str(target))

        mcp, _ = _setup(make_config(), modules=[batch_structured])
        out = _tool(mcp, "batch_replace_structured")(
            glob="documents:/deck.pptx",
            old="hit me",
            new="X",
            dry_run=False,
        )
        assert out["count"] == 1

    def test_unsupported_format_propagates(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        """Inner replace helper raising :class:`UnsupportedFormatError` must
        bubble out without being swallowed (line 209-210)."""
        from docx import Document

        docs_dir, _ = tmp_roots
        target = docs_dir / "x.docx"
        Document().save(str(target))
        mcp, _ = _setup(make_config(), modules=[batch_structured])
        with patch.object(
            batch_structured,
            "_replace_in_docx",
            side_effect=UnsupportedFormatError("nope"),
        ):
            with pytest.raises(UnsupportedFormatError):
                _tool(mcp, "batch_replace_structured")(
                    glob="documents:/*.docx",
                    old="x",
                    new="y",
                    dry_run=True,
                )

    def test_pre_snapshot_failure_is_silent(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        """Snapshot failures must not block the apply pass (line 240-241)."""
        from docx import Document

        docs_dir, _ = tmp_roots
        target = docs_dir / "x.docx"
        d = Document()
        d.add_paragraph("hit me")
        d.save(str(target))
        mcp, ctx = _setup(make_config(), modules=[batch_structured])

        original = ctx.versions.snapshot
        calls: list[str] = []

        def flaky(*args, **kwargs):  # type: ignore[no-untyped-def]
            action = kwargs.get("action", "")
            calls.append(action)
            if action.endswith("_pre"):
                raise RuntimeError("pre snapshot blew up")
            return original(*args, **kwargs)

        with patch.object(ctx.versions, "snapshot", side_effect=flaky):
            out = _tool(mcp, "batch_replace_structured")(
                glob="documents:/*.docx",
                old="hit me",
                new="X",
                dry_run=False,
            )
        assert out["count"] == 1
        assert any(c.endswith("_pre") for c in calls)

    def test_pptx_paragraph_no_match_continues(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        """Cover the n==0 fall-through (branch 139→134) — a paragraph with
        text but no pattern match must skip to the next paragraph."""
        from pptx import Presentation
        from pptx.util import Inches

        docs_dir, _ = tmp_roots
        target = docs_dir / "deck.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(3))
        tf = tx.text_frame
        tf.text = "no-match here"
        tf.add_paragraph().text = "hit me"
        prs.save(str(target))

        mcp, _ = _setup(make_config(), modules=[batch_structured])
        out = _tool(mcp, "batch_replace_structured")(
            glob="documents:/deck.pptx",
            old="hit me",
            new="X",
            dry_run=False,
        )
        assert out["count"] == 1

    def test_pptx_paragraph_multiple_runs_cleanup(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        """Cover line 153 — when a paragraph has >1 run, the loop clears
        the tail runs before writing into runs[0]."""
        from pptx import Presentation
        from pptx.util import Inches

        docs_dir, _ = tmp_roots
        target = docs_dir / "multi.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(3))
        tf = tx.text_frame
        # Build a paragraph with three distinct runs that concatenate to
        # "hit me here".
        para = tf.paragraphs[0]
        para.text = "hit "
        run2 = para.add_run()
        run2.text = "me "
        run3 = para.add_run()
        run3.text = "here"
        prs.save(str(target))

        mcp, _ = _setup(make_config(), modules=[batch_structured])
        out = _tool(mcp, "batch_replace_structured")(
            glob="documents:/multi.pptx",
            old="hit me here",
            new="REPLACED",
            dry_run=False,
        )
        assert out["count"] == 1
        # Verify the replacement actually happened in the file.
        prs2 = Presentation(str(target))
        all_text = "\n".join(
            sh.text_frame.text
            for slide in prs2.slides
            for sh in slide.shapes
            if hasattr(sh, "text_frame")
        )
        assert "REPLACED" in all_text

    def test_post_snapshot_failure_is_silent(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        from docx import Document

        docs_dir, _ = tmp_roots
        target = docs_dir / "x.docx"
        d = Document()
        d.add_paragraph("hit me")
        d.save(str(target))
        mcp, ctx = _setup(make_config(), modules=[batch_structured])

        original = ctx.versions.snapshot

        def flaky(*args, **kwargs):  # type: ignore[no-untyped-def]
            if kwargs.get("action", "").endswith("_post"):
                raise RuntimeError("post snapshot blew up")
            return original(*args, **kwargs)

        with patch.object(ctx.versions, "snapshot", side_effect=flaky):
            out = _tool(mcp, "batch_replace_structured")(
                glob="documents:/*.docx",
                old="hit me",
                new="X",
                dry_run=False,
            )
        assert out["count"] == 1
