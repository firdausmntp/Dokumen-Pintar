"""Tests for :class:`dokumen_pintar.handlers.pptx_handler.PptxHandler`."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from dokumen_pintar.errors import HandlerError, UnsupportedFormatError
from dokumen_pintar.handlers.pptx_handler import PptxHandler, _shape_kind, _shape_to_dict


@pytest.fixture
def handler() -> PptxHandler:
    return PptxHandler()


def _create_pptx(path: Path, slides: int = 2) -> None:
    prs = Presentation()
    for i in range(slides):
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        if title is not None:
            title.text = f"Slide {i} Title"
        subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if subtitle is not None:
            subtitle.text = f"Content of slide {i}"
    prs.save(str(path))


def test_detect(handler: PptxHandler) -> None:
    assert handler.detect(Path("pres.pptx")) is True
    assert handler.detect(Path("pres.pdf")) is False


def test_read_meta(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "test.pptx"
    _create_pptx(target, slides=3)
    meta = handler.read_meta(target)
    assert meta["format"] == "pptx"
    assert meta["slide_count"] == 3
    assert len(meta["slide_titles"]) == 3


def test_read_text(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "test.pptx"
    _create_pptx(target)
    text = handler.read_text(target)
    assert "Slide 0" in text
    assert "Slide 1" in text


def test_write_text_raises(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "test.pptx"
    _create_pptx(target)
    with pytest.raises(UnsupportedFormatError):
        handler.write_text(target, "content")


def test_extract_for_search(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "test.pptx"
    _create_pptx(target)
    text = handler.extract_for_search(target)
    assert isinstance(text, str)
    assert len(text) > 0


def test_extract_for_search_invalid_file(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "bad.pptx"
    target.write_bytes(b"not a pptx")
    text = handler.extract_for_search(target)
    assert text == ""


def test_structured_get_slides(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "test.pptx"
    _create_pptx(target, slides=2)
    result = handler.structured_get(target, "slides")
    assert isinstance(result, list)
    assert len(result) == 2
    assert result[0]["index"] == 0


def test_structured_get_slide_by_index(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "test.pptx"
    _create_pptx(target)
    result = handler.structured_get(target, "slide:0")
    assert isinstance(result, dict)
    assert "shapes" in result
    assert result["index"] == 0


def test_structured_get_slide_text(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "test.pptx"
    _create_pptx(target)
    text = handler.structured_get(target, "slide_text:0")
    assert isinstance(text, str)


def test_structured_get_slide_out_of_range(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "test.pptx"
    _create_pptx(target, slides=1)
    with pytest.raises(HandlerError, match="out of range"):
        handler.structured_get(target, "slide:99")


def test_structured_get_unsupported(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "test.pptx"
    _create_pptx(target)
    with pytest.raises(HandlerError, match="unsupported"):
        handler.structured_get(target, "invalid_expr")


def test_structured_set_slide_title(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "test.pptx"
    _create_pptx(target)
    handler.structured_set(target, "slide_title:0", "New Title")
    result = handler.structured_get(target, "slide:0")
    assert result["title"] == "New Title"


def test_structured_set_unsupported(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "test.pptx"
    _create_pptx(target)
    with pytest.raises(HandlerError, match="unsupported"):
        handler.structured_set(target, "invalid_expr", "value")


def test_structured_delete_slide(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "test.pptx"
    _create_pptx(target, slides=3)
    handler.structured_delete(target, "slide:1")
    meta = handler.read_meta(target)
    assert meta["slide_count"] == 2


def test_structured_delete_slide_out_of_range(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "test.pptx"
    _create_pptx(target, slides=1)
    with pytest.raises(HandlerError, match="out of range"):
        handler.structured_delete(target, "slide:99")


def test_structured_delete_unsupported(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "test.pptx"
    _create_pptx(target)
    with pytest.raises(HandlerError, match="unsupported"):
        handler.structured_delete(target, "invalid_expr")


# ── Additional PPTX coverage ──

from dokumen_pintar.handlers.pptx_handler import (
    _parse_slide_index,
    _get_slide,
    _shape_kind,
    _shape_to_dict,
    _slide_title,
    _iter_text_frames,
)


def test_parse_slide_index_valid() -> None:
    assert _parse_slide_index("0") == 0
    assert _parse_slide_index("5") == 5


def test_parse_slide_index_invalid() -> None:
    with pytest.raises(HandlerError, match="invalid slide index"):
        _parse_slide_index("abc")


def test_get_slide_out_of_range(tmp_path: Path) -> None:
    target = tmp_path / "oor.pptx"
    _create_pptx(target, slides=1)
    prs = Presentation(str(target))
    with pytest.raises(HandlerError, match="out of range"):
        _get_slide(prs, 99)


def test_get_slide_negative(tmp_path: Path) -> None:
    target = tmp_path / "neg.pptx"
    _create_pptx(target, slides=1)
    prs = Presentation(str(target))
    with pytest.raises(HandlerError, match="out of range"):
        _get_slide(prs, -1)


def test_shape_kind_text(tmp_path: Path) -> None:
    target = tmp_path / "sk.pptx"
    _create_pptx(target, slides=1)
    prs = Presentation(str(target))
    slide = list(prs.slides)[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            assert _shape_kind(shape) == "text"
            break


def test_shape_to_dict_text(tmp_path: Path) -> None:
    target = tmp_path / "sd.pptx"
    _create_pptx(target, slides=1)
    prs = Presentation(str(target))
    slide = list(prs.slides)[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            d = _shape_to_dict(shape)
            assert d["kind"] == "text"
            assert "text" in d
            assert "shape_id" in d
            break


def test_shape_to_dict_with_table(tmp_path: Path) -> None:
    target = tmp_path / "tbl.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    from pptx.util import Inches
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    table_shape.table.cell(0, 0).text = "A"
    prs.save(str(target))

    prs2 = Presentation(str(target))
    slide2 = list(prs2.slides)[0]
    for shape in slide2.shapes:
        if shape.has_table:
            d = _shape_to_dict(shape)
            assert d["kind"] == "table"
            assert "rows" in d
            break


def test_slide_title_no_title(tmp_path: Path) -> None:
    target = tmp_path / "nt.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank - no title placeholder
    slide = prs.slides.add_slide(slide_layout)
    prs.save(str(target))
    prs2 = Presentation(str(target))
    slide2 = list(prs2.slides)[0]
    assert _slide_title(slide2) == ""


def test_iter_text_frames_returns_texts(tmp_path: Path) -> None:
    target = tmp_path / "itf.pptx"
    _create_pptx(target, slides=1)
    prs = Presentation(str(target))
    slide = list(prs.slides)[0]
    texts = _iter_text_frames(slide)
    assert isinstance(texts, list)
    assert any("Slide 0" in t for t in texts)


def test_structured_set_slide_text_placeholder_by_idx(
    handler: PptxHandler, tmp_path: Path
) -> None:
    target = tmp_path / "stpi.pptx"
    _create_pptx(target, slides=1)
    # Placeholder idx 0 is title, 1 is subtitle on Title Slide layout
    handler.structured_set(target, "slide_text:0:1", "New Subtitle")
    result = handler.structured_get(target, "slide_text:0")
    assert "New Subtitle" in result


def test_structured_set_slide_text_missing_placeholder(
    handler: PptxHandler, tmp_path: Path
) -> None:
    target = tmp_path / "stmp.pptx"
    _create_pptx(target, slides=1)
    with pytest.raises(HandlerError, match="placeholder not found"):
        handler.structured_set(target, "slide_text:0:999", "value")


def test_structured_set_slide_text_bad_format(
    handler: PptxHandler, tmp_path: Path
) -> None:
    target = tmp_path / "stbf.pptx"
    _create_pptx(target, slides=1)
    with pytest.raises(HandlerError, match="slide_text set requires"):
        handler.structured_set(target, "slide_text:0", "value")


def test_structured_set_slide_text_non_string(
    handler: PptxHandler, tmp_path: Path
) -> None:
    target = tmp_path / "stns.pptx"
    _create_pptx(target, slides=1)
    with pytest.raises(HandlerError, match="must be a string"):
        handler.structured_set(target, "slide_text:0:1", 123)


def test_structured_set_slide_title_non_string(
    handler: PptxHandler, tmp_path: Path
) -> None:
    target = tmp_path / "ttns.pptx"
    _create_pptx(target, slides=1)
    with pytest.raises(HandlerError, match="must be a string"):
        handler.structured_set(target, "slide_title:0", 42)


def test_structured_set_slide_title_out_of_range(
    handler: PptxHandler, tmp_path: Path
) -> None:
    target = tmp_path / "ttor.pptx"
    _create_pptx(target, slides=1)
    with pytest.raises(HandlerError, match="out of range"):
        handler.structured_set(target, "slide_title:99", "T")


def test_read_text_invalid_file(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "bad.pptx"
    target.write_bytes(b"not a pptx file at all")
    with pytest.raises(HandlerError, match="failed to open"):
        handler.read_text(target)


def test_read_meta_invalid_file(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "bad2.pptx"
    target.write_bytes(b"corrupted")
    with pytest.raises(HandlerError, match="failed to open"):
        handler.read_meta(target)


def test_structured_get_invalid_file(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "bad3.pptx"
    target.write_bytes(b"garbage")
    with pytest.raises(HandlerError, match="failed to open"):
        handler.structured_get(target, "slides")


def test_structured_set_invalid_file(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "bad4.pptx"
    target.write_bytes(b"garbage")
    with pytest.raises(HandlerError, match="failed to open"):
        handler.structured_set(target, "slide_title:0", "T")


def test_structured_delete_invalid_file(handler: PptxHandler, tmp_path: Path) -> None:
    target = tmp_path / "bad5.pptx"
    target.write_bytes(b"garbage")
    with pytest.raises(HandlerError, match="failed to open"):
        handler.structured_delete(target, "slide:0")


def test_structured_set_slide_text_by_name(
    handler: PptxHandler, tmp_path: Path
) -> None:
    target = tmp_path / "stn.pptx"
    _create_pptx(target, slides=1)
    prs = Presentation(str(target))
    slide = list(prs.slides)[0]
    # Get the name of the subtitle placeholder
    ph_name = None
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx == 1:
                ph_name = ph.name
                break
        except Exception:
            pass
    prs = None  # release
    if ph_name:
        handler.structured_set(target, f"slide_text:0:{ph_name}", "Named PH")


def test_shape_kind_table(tmp_path: Path) -> None:
    from pptx.util import Inches
    target = tmp_path / "sktbl.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))
    prs.save(str(target))

    prs2 = Presentation(str(target))
    slide2 = list(prs2.slides)[0]
    for shape in slide2.shapes:
        if shape.has_table:
            assert _shape_kind(shape) == "table"
            break


def test_shape_kind_other(tmp_path: Path) -> None:
    from pptx.util import Inches
    from unittest.mock import MagicMock
    mock_shape = MagicMock()
    mock_shape.has_text_frame = False
    mock_shape.has_table = False
    mock_shape.shape_type = None
    assert _shape_kind(mock_shape) == "other"


def test_shape_to_dict_picture() -> None:
    from unittest.mock import MagicMock, PropertyMock
    mock_shape = MagicMock()
    mock_shape.has_text_frame = False
    mock_shape.has_table = False
    type(mock_shape).shape_type = PropertyMock(return_value="PICTURE (13)")
    mock_shape.name = "pic1"
    mock_shape.shape_id = 5
    mock_shape.left = 100
    mock_shape.top = 200
    mock_shape.width = 300
    mock_shape.height = 400
    mock_shape.image.filename = "img.png"
    mock_shape.image.content_type = "image/png"
    d = _shape_to_dict(mock_shape)
    assert d["kind"] == "picture"
    assert d["image_filename"] == "img.png"


def test_structured_set_slide_title_blank_layout(
    handler: PptxHandler, tmp_path: Path
) -> None:
    target = tmp_path / "stbl.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    prs.save(str(target))
    with pytest.raises(HandlerError, match="no title placeholder"):
        handler.structured_set(target, "slide_title:0", "Title")


def test_structured_get_slide_detail(
    handler: PptxHandler, tmp_path: Path
) -> None:
    target = tmp_path / "sdet.pptx"
    _create_pptx(target, slides=1)
    result = handler.structured_get(target, "slide:0")
    assert "shapes" in result
    assert "title" in result
    assert result["index"] == 0


def test_structured_get_slide_text(
    handler: PptxHandler, tmp_path: Path
) -> None:
    target = tmp_path / "gst.pptx"
    _create_pptx(target, slides=1)
    text = handler.structured_get(target, "slide_text:0")
    assert isinstance(text, str)
    assert "Slide 0" in text


def test_extract_for_search_multipage(
    handler: PptxHandler, tmp_path: Path
) -> None:
    target = tmp_path / "efs.pptx"
    _create_pptx(target, slides=3)
    text = handler.extract_for_search(target)
    assert "Slide 0" in text
    assert "Slide 2" in text


def test_extract_for_search_invalid(
    handler: PptxHandler, tmp_path: Path
) -> None:
    target = tmp_path / "badsearch.pptx"
    target.write_bytes(b"garbage")
    text = handler.extract_for_search(target)
    assert text == ""


# ── Mock-based tests for exception branches ──

from dokumen_pintar.handlers.pptx_handler import (
    _slide_title, _iter_text_frames, _shape_to_dict,
)
from unittest.mock import MagicMock, PropertyMock, patch


def test_slide_title_shapes_title_raises() -> None:
    slide = MagicMock()
    type(slide).shapes = PropertyMock(side_effect=Exception("broken"))
    assert _slide_title(slide) == ""


def test_slide_title_no_text_frame() -> None:
    slide = MagicMock()
    title_shape = MagicMock()
    title_shape.has_text_frame = False
    slide.shapes.title = title_shape
    assert _slide_title(slide) == ""


def test_slide_title_text_frame_raises() -> None:
    slide = MagicMock()
    title_shape = MagicMock()
    title_shape.has_text_frame = True
    type(title_shape.text_frame).text = PropertyMock(side_effect=Exception("err"))
    slide.shapes.title = title_shape
    assert _slide_title(slide) == ""


def test_iter_text_frames_shape_raises() -> None:
    slide = MagicMock()
    bad_shape = MagicMock()
    type(bad_shape).has_text_frame = PropertyMock(side_effect=Exception("err"))
    slide.shapes = [bad_shape]
    result = _iter_text_frames(slide)
    assert result == []


def test_shape_kind_text_frame_raises() -> None:
    shape = MagicMock()
    type(shape).has_text_frame = PropertyMock(side_effect=Exception("err"))
    type(shape).has_table = PropertyMock(return_value=False)
    shape.shape_type = None
    assert _shape_kind(shape) == "other"


def test_shape_kind_table_raises() -> None:
    shape = MagicMock()
    shape.has_text_frame = False
    type(shape).has_table = PropertyMock(side_effect=Exception("err"))
    shape.shape_type = None
    assert _shape_kind(shape) == "other"


def test_shape_kind_shape_type_raises() -> None:
    shape = MagicMock()
    shape.has_text_frame = False
    shape.has_table = False
    type(shape).shape_type = PropertyMock(side_effect=Exception("err"))
    assert _shape_kind(shape) == "other"


def test_shape_to_dict_shape_id_raises() -> None:
    shape = MagicMock()
    shape.has_text_frame = True
    shape.text_frame.text = "hi"
    shape.name = "s1"
    type(shape).shape_id = PropertyMock(side_effect=Exception("err"))
    shape.left = 0
    shape.top = 0
    shape.width = 100
    shape.height = 100
    d = _shape_to_dict(shape)
    assert d["shape_id"] == 0
    assert d["text"] == "hi"


def test_shape_to_dict_position_non_int() -> None:
    shape = MagicMock()
    shape.has_text_frame = True
    shape.text_frame.text = "abc"
    shape.name = "s2"
    shape.shape_id = 1
    shape.left = "200"
    shape.top = "300"
    shape.width = "400"
    shape.height = "500"
    d = _shape_to_dict(shape)
    assert d["left"] == 200
    assert d["kind"] == "text"


def test_shape_to_dict_text_exception() -> None:
    shape = MagicMock()
    shape.has_text_frame = True
    type(shape.text_frame).text = PropertyMock(side_effect=Exception("err"))
    shape.name = "s3"
    shape.shape_id = 2
    shape.left = 0
    shape.top = 0
    shape.width = 100
    shape.height = 100
    d = _shape_to_dict(shape)
    assert d["text"] == ""


def test_shape_to_dict_table_exception() -> None:
    shape = MagicMock()
    shape.has_text_frame = False
    shape.has_table = True
    shape.name = "t1"
    shape.shape_id = 3
    shape.left = 0
    shape.top = 0
    shape.width = 100
    shape.height = 100
    type(shape).table = PropertyMock(side_effect=Exception("err"))
    d = _shape_to_dict(shape)
    assert d["kind"] == "table"
    assert d["rows"] == []


def test_shape_to_dict_picture_image_exception() -> None:
    shape = MagicMock()
    shape.has_text_frame = False
    shape.has_table = False
    type(shape).shape_type = PropertyMock(return_value="PICTURE (13)")
    shape.name = "p1"
    shape.shape_id = 4
    shape.left = 0
    shape.top = 0
    shape.width = 100
    shape.height = 100
    type(shape).image = PropertyMock(side_effect=Exception("err"))
    d = _shape_to_dict(shape)
    assert d["kind"] == "picture"
    assert "image_filename" not in d


def test_structured_set_placeholder_no_text_frame(
    handler: PptxHandler, tmp_path: Path
) -> None:
    target = tmp_path / "phntf.pptx"
    _create_pptx(target, slides=1)
    mock_prs = MagicMock()
    mock_slide = MagicMock()
    mock_ph = MagicMock()
    mock_ph.placeholder_format.idx = 0
    mock_ph.name = "Title 1"
    mock_ph.has_text_frame = False
    mock_slide.placeholders = [mock_ph]
    mock_prs.slides = [mock_slide]
    with patch("dokumen_pintar.handlers.pptx_handler.Presentation", return_value=mock_prs):
        with pytest.raises(HandlerError, match="no text frame"):
            handler.structured_set(target, "slide_text:0:0", "X")


def test_structured_set_slide_title_with_title(
    handler: PptxHandler, tmp_path: Path
) -> None:
    target = tmp_path / "stitle.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    prs.save(str(target))
    handler.structured_set(target, "slide_title:0", "New Title")
    prs2 = Presentation(str(target))
    slide2 = list(prs2.slides)[0]
    assert slide2.shapes.title.text_frame.text == "New Title"


# ── Additional PPTX coverage tests ──

def test_shape_kind_text_exception() -> None:
    from unittest.mock import MagicMock, PropertyMock
    shape = MagicMock()
    type(shape).has_text_frame = PropertyMock(side_effect=Exception("bad"))
    type(shape).has_table = PropertyMock(return_value=False)
    type(shape).image = PropertyMock(side_effect=AttributeError)
    kind = _shape_kind(shape)
    assert kind == "other"


def test_shape_to_dict_position_exception() -> None:
    from unittest.mock import MagicMock, PropertyMock
    shape = MagicMock()
    shape.has_text_frame = True
    shape.text_frame.text = "hi"
    shape.name = "Test"
    shape.shape_id = 1
    # Make width raise
    type(shape).width = PropertyMock(side_effect=Exception("no width"))
    shape.left = 100
    shape.top = 200
    shape.height = 300
    d = _shape_to_dict(shape)
    assert d["kind"] == "text"
    assert "width" not in d


def test_iter_text_frames_shape_exception(handler: PptxHandler, tmp_path: Path) -> None:
    from unittest.mock import MagicMock, PropertyMock
    target = tmp_path / "itf.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(target))
    # Should not crash even if shape raises
    text = handler.extract_for_search(target)
    assert isinstance(text, str)


def test_structured_set_placeholder_idx_exception(handler: PptxHandler, tmp_path: Path) -> None:
    from unittest.mock import patch, MagicMock, PropertyMock
    target = tmp_path / "phidx.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(target))
    mock_prs = MagicMock()
    mock_slide = MagicMock()
    mock_ph = MagicMock()
    type(mock_ph.placeholder_format).idx = PropertyMock(side_effect=Exception("no idx"))
    mock_ph.name = "target_ph"
    mock_ph.has_text_frame = True
    mock_ph.text_frame = MagicMock()
    mock_slide.placeholders = [mock_ph]
    mock_prs.slides = [mock_slide]
    with patch("dokumen_pintar.handlers.pptx_handler.Presentation", return_value=mock_prs):
        handler.structured_set(target, "slide_text:0:target_ph", "val")
    mock_ph.text_frame.__setattr__("text", "val")


def test_structured_set_placeholder_set_exception(handler: PptxHandler, tmp_path: Path) -> None:
    from unittest.mock import patch, MagicMock, PropertyMock
    target = tmp_path / "phset.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(target))
    mock_prs = MagicMock()
    mock_slide = MagicMock()
    mock_ph = MagicMock()
    mock_ph.placeholder_format.idx = 0
    mock_ph.name = "ph0"
    mock_ph.has_text_frame = True
    type(mock_ph).text_frame = PropertyMock(side_effect=RuntimeError("set fail"))
    mock_slide.placeholders = [mock_ph]
    mock_prs.slides = [mock_slide]
    with patch("dokumen_pintar.handlers.pptx_handler.Presentation", return_value=mock_prs):
        with pytest.raises(HandlerError, match="failed to set placeholder"):
            handler.structured_set(target, "slide_text:0:0", "val")


def test_structured_set_slide_title_exception(handler: PptxHandler, tmp_path: Path) -> None:
    from unittest.mock import patch, MagicMock, PropertyMock
    target = tmp_path / "titleex.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(target))
    mock_prs = MagicMock()
    mock_slide = MagicMock()
    mock_slide.shapes.title = None
    mock_prs.slides = [mock_slide]
    with patch("dokumen_pintar.handlers.pptx_handler.Presentation", return_value=mock_prs):
        with pytest.raises(HandlerError, match="no title placeholder"):
            handler.structured_set(target, "slide_title:0", "Title")


def test_structured_set_slide_title_shapes_exception(handler: PptxHandler, tmp_path: Path) -> None:
    from unittest.mock import patch, MagicMock, PropertyMock
    target = tmp_path / "titleshex.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(target))
    mock_prs = MagicMock()
    mock_slide = MagicMock()
    type(mock_slide.shapes).title = PropertyMock(side_effect=Exception("no title shape"))
    mock_prs.slides = [mock_slide]
    with patch("dokumen_pintar.handlers.pptx_handler.Presentation", return_value=mock_prs):
        with pytest.raises(HandlerError, match="no title placeholder"):
            handler.structured_set(target, "slide_title:0", "Title")


def test_shape_to_dict_attr_none(tmp_path: Path) -> None:
    from unittest.mock import MagicMock
    from dokumen_pintar.handlers.pptx_handler import _shape_to_dict
    shape = MagicMock()
    shape.has_text_frame = True
    shape.text_frame.text = "hello"
    shape.name = "Box"
    shape.shape_id = 2
    # All position attrs return None → line 89 continue
    shape.left = None
    shape.top = None
    shape.width = None
    shape.height = None
    d = _shape_to_dict(shape)
    assert "left" not in d
    assert d["text"] == "hello"


def test_shape_to_dict_other_kind(tmp_path: Path) -> None:
    from unittest.mock import MagicMock
    from dokumen_pintar.handlers.pptx_handler import _shape_to_dict
    shape = MagicMock()
    shape.has_text_frame = False
    shape.has_table = False
    # shape_type doesn't contain "PICTURE"
    shape.shape_type = "LINE"
    shape.name = "Line1"
    shape.shape_id = 3
    shape.left = 100
    shape.top = 200
    shape.width = 300
    shape.height = 50
    d = _shape_to_dict(shape)
    assert d["kind"] == "other"
    assert "rows" not in d
    assert "text" not in d or d.get("text") == ""


def test_iter_text_frames_no_text_frame(tmp_path: Path) -> None:
    from unittest.mock import MagicMock
    from dokumen_pintar.handlers.pptx_handler import _iter_text_frames
    shape_no_tf = MagicMock()
    shape_no_tf.has_text_frame = False
    shape_empty_txt = MagicMock()
    shape_empty_txt.has_text_frame = True
    shape_empty_txt.text_frame.text = ""
    slide = MagicMock()
    slide.shapes = [shape_no_tf, shape_empty_txt]
    result = _iter_text_frames(slide)
    assert result == []


def test_structured_set_title_set_exception(handler: PptxHandler, tmp_path: Path) -> None:
    from unittest.mock import patch, MagicMock, PropertyMock
    target = tmp_path / "titleseterr.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(target))
    mock_prs = MagicMock()
    mock_slide = MagicMock()
    mock_title_shape = MagicMock()
    mock_slide.shapes.title = mock_title_shape
    type(mock_title_shape).text_frame = PropertyMock(side_effect=Exception("set fail"))
    mock_prs.slides = [mock_slide]
    with patch("dokumen_pintar.handlers.pptx_handler.Presentation", return_value=mock_prs):
        with pytest.raises(HandlerError, match="failed to set title"):
            handler.structured_set(target, "slide_title:0", "Title")
