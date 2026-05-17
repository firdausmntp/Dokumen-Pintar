"""Tests for :mod:`dokumen_pintar.tools.images`."""

from __future__ import annotations

import io
import zipfile
from pathlib import Path
from typing import Callable

import pytest
from docx import Document
from docx.shared import Inches
from mcp.server.fastmcp import FastMCP
from PIL import Image as PilImage
from pptx import Presentation
from pptx.util import Inches as PptxInches

from dokumen_pintar.config import AppConfig
from dokumen_pintar.context import build_context
from dokumen_pintar.errors import UnsupportedFormatError, ValidationError
from dokumen_pintar.tools import images


def _setup(cfg: AppConfig) -> tuple[FastMCP, ...]:
    ctx = build_context(cfg)
    mcp = FastMCP(name="t-images")
    images.register(mcp, ctx)
    return mcp, ctx


def _tool(mcp: FastMCP, name: str):
    return mcp._tool_manager._tools[name].fn


def _png_bytes(color: str = "red", size: tuple[int, int] = (32, 32)) -> bytes:
    buf = io.BytesIO()
    PilImage.new("RGB", size, color).save(buf, format="PNG")
    return buf.getvalue()


def _make_docx_with_images(target: Path, n_images: int = 2) -> None:
    doc = Document()
    doc.add_paragraph("hello")
    # python-docx dedupes images that hash to the same bytes - use unique
    # colors AND sizes per image so each lands as a distinct media entry.
    palette = ["red", "blue", "green", "yellow", "purple", "orange"]
    for i in range(n_images):
        img_buf = io.BytesIO(_png_bytes(palette[i % len(palette)], size=(32 + i, 32 + i)))
        doc.add_picture(img_buf, width=Inches(1))
    doc.save(str(target))


def _make_pptx_with_image(target: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    img_buf = io.BytesIO(_png_bytes("green"))
    slide.shapes.add_picture(img_buf, PptxInches(1), PptxInches(1), PptxInches(2), PptxInches(2))
    prs.save(target)


# ── image_list ──


def test_image_list_docx(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    p = docs_dir / "with_imgs.docx"
    _make_docx_with_images(p, n_images=3)
    mcp, _ctx = _setup(make_config())
    result = _tool(mcp, "image_list")(path=str(p))
    assert result["format"] == "docx"
    assert result["count"] == 3
    assert all(entry["ext"] == ".png" for entry in result["images"])
    assert all(entry["size"] > 0 for entry in result["images"])


def test_image_list_pptx(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    p = docs_dir / "deck.pptx"
    _make_pptx_with_image(p)
    mcp, _ctx = _setup(make_config())
    result = _tool(mcp, "image_list")(path=str(p))
    assert result["format"] == "pptx"
    assert result["count"] >= 1


def test_image_list_unsupported_format_raises(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    p = docs_dir / "data.xlsx"
    import openpyxl

    openpyxl.Workbook().save(p)
    mcp, _ctx = _setup(make_config())
    with pytest.raises(UnsupportedFormatError, match="image_list"):
        _tool(mcp, "image_list")(path=str(p))


def test_image_list_no_handler_raises(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    p = docs_dir / "weird.unknownext"
    p.write_text("?", encoding="utf-8")
    mcp, _ctx = _setup(make_config())
    with pytest.raises(UnsupportedFormatError, match="no handler"):
        _tool(mcp, "image_list")(path=str(p))


def test_image_list_corrupted_zip_raises(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    from dokumen_pintar.errors import HandlerError

    docs_dir, _ = tmp_roots
    p = docs_dir / "broken.docx"
    p.write_bytes(b"PK\x03\x04not-a-real-zip")
    mcp, _ctx = _setup(make_config())
    with pytest.raises(HandlerError, match="not a valid zip"):
        _tool(mcp, "image_list")(path=str(p))


# ── image_extract ──


def test_image_extract_docx_writes_blob(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    src = docs_dir / "src.docx"
    _make_docx_with_images(src, n_images=2)
    mcp, _ctx = _setup(make_config())
    out = docs_dir / "extracted_0.png"
    result = _tool(mcp, "image_extract")(path=str(src), index=0, dst=str(out))
    assert Path(result["dst"]).exists()
    assert result["size"] > 0
    assert result["ext"] == ".png"
    # File on disk is a valid PNG.
    PilImage.open(result["dst"]).verify()


def test_image_extract_corrects_extension(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    """Destination ext is forced to match the source image's actual ext."""
    docs_dir, _ = tmp_roots
    src = docs_dir / "src2.docx"
    _make_docx_with_images(src, n_images=1)
    mcp, _ctx = _setup(make_config())
    # Caller passed .jpg but source is PNG - tool must rename to .png.
    out = docs_dir / "weird.jpg"
    result = _tool(mcp, "image_extract")(path=str(src), index=0, dst=str(out))
    assert result["dst"].endswith(".png")
    assert not Path(str(out)).exists()  # original .jpg name not used


def test_image_extract_rejects_overwrite(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    src = docs_dir / "src3.docx"
    _make_docx_with_images(src, n_images=1)
    out = docs_dir / "exists.png"
    out.write_bytes(b"old")
    mcp, _ctx = _setup(make_config())
    with pytest.raises(ValidationError, match="destination exists"):
        _tool(mcp, "image_extract")(path=str(src), index=0, dst=str(out))
    # overwrite=True works.
    _tool(mcp, "image_extract")(path=str(src), index=0, dst=str(out), overwrite=True)
    assert out.read_bytes() != b"old"


def test_image_extract_out_of_range(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    src = docs_dir / "src4.docx"
    _make_docx_with_images(src, n_images=1)
    out = docs_dir / "nope.png"
    mcp, _ctx = _setup(make_config())
    with pytest.raises(ValidationError, match="out of range"):
        _tool(mcp, "image_extract")(path=str(src), index=99, dst=str(out))


def test_image_extract_negative_index(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    src = docs_dir / "src5.docx"
    _make_docx_with_images(src, n_images=1)
    mcp, _ctx = _setup(make_config())
    with pytest.raises(ValidationError, match=">= 0"):
        _tool(mcp, "image_extract")(path=str(src), index=-1, dst="x.png")


def test_image_extract_unsupported_format(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    """xlsx has no images to extract - tool refuses."""
    import openpyxl

    docs_dir, _ = tmp_roots
    src = docs_dir / "data.xlsx"
    openpyxl.Workbook().save(src)
    out = docs_dir / "out.png"
    mcp, _ctx = _setup(make_config())
    with pytest.raises(UnsupportedFormatError, match="image_extract"):
        _tool(mcp, "image_extract")(path=str(src), index=0, dst=str(out))


# ── image_extract_all ──


def test_image_extract_all_writes_every_image(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    src = docs_dir / "all.docx"
    _make_docx_with_images(src, n_images=3)
    out_dir = docs_dir / "imgs_out"
    mcp, _ctx = _setup(make_config())
    result = _tool(mcp, "image_extract_all")(path=str(src), dst_dir=str(out_dir))
    assert result["count"] == 3
    written = list(out_dir.iterdir())
    assert len(written) == 3
    for f in written:
        assert f.suffix == ".png"


def test_image_extract_all_custom_naming(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    src = docs_dir / "named.docx"
    _make_docx_with_images(src, n_images=2)
    out_dir = docs_dir / "named_out"
    mcp, _ctx = _setup(make_config())
    result = _tool(mcp, "image_extract_all")(
        path=str(src),
        dst_dir=str(out_dir),
        naming_pattern="img-{index}{ext}",
    )
    names = sorted(Path(f["path"]).name for f in result["files"])
    assert names == ["img-0.png", "img-1.png"]


def test_image_extract_all_unsupported_format(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    import openpyxl

    docs_dir, _ = tmp_roots
    src = docs_dir / "ea.xlsx"
    openpyxl.Workbook().save(src)
    mcp, _ctx = _setup(make_config())
    with pytest.raises(UnsupportedFormatError, match="image_extract_all"):
        _tool(mcp, "image_extract_all")(path=str(src), dst_dir=str(docs_dir / "dummy"))


# ── image_replace ──


def test_image_replace_swaps_bytes(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    target = docs_dir / "rep.docx"
    _make_docx_with_images(target, n_images=2)

    # Create a fresh standalone PNG to inject.
    new_png = docs_dir / "new.png"
    new_png.write_bytes(_png_bytes("yellow", size=(64, 64)))

    mcp, _ctx = _setup(make_config())
    listed = _tool(mcp, "image_list")(path=str(target))
    old_size = listed["images"][0]["size"]
    new_size = new_png.stat().st_size
    assert old_size != new_size  # sanity

    result = _tool(mcp, "image_replace")(path=str(target), index=0, src=str(new_png))
    assert result["size_old"] == old_size
    assert result["size_new"] == new_size

    # Re-listing reflects the new size at index 0.
    re_listed = _tool(mcp, "image_list")(path=str(target))
    assert re_listed["images"][0]["size"] == new_size
    # Internal name preserved (existing references in document still work).
    assert re_listed["images"][0]["internal_name"] == listed["images"][0]["internal_name"]


def test_image_replace_unsupported_pdf(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    """PDF is read-only - image_replace must refuse."""
    from reportlab.pdfgen import canvas

    docs_dir, _ = tmp_roots
    pdf_path = docs_dir / "stub.pdf"
    c = canvas.Canvas(str(pdf_path))
    c.drawString(100, 750, "stub")
    c.save()
    new_png = docs_dir / "x.png"
    new_png.write_bytes(_png_bytes())
    mcp, _ctx = _setup(make_config())
    with pytest.raises(UnsupportedFormatError, match="image_replace"):
        _tool(mcp, "image_replace")(path=str(pdf_path), index=0, src=str(new_png))


def test_image_replace_out_of_range(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    target = docs_dir / "oor.docx"
    _make_docx_with_images(target, n_images=1)
    new_png = docs_dir / "y.png"
    new_png.write_bytes(_png_bytes())
    mcp, _ctx = _setup(make_config())
    with pytest.raises(ValidationError, match="out of range"):
        _tool(mcp, "image_replace")(path=str(target), index=99, src=str(new_png))


def test_image_replace_negative_index(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    target = docs_dir / "neg.docx"
    _make_docx_with_images(target, n_images=1)
    new_png = docs_dir / "z.png"
    new_png.write_bytes(_png_bytes())
    mcp, _ctx = _setup(make_config())
    with pytest.raises(ValidationError, match=">= 0"):
        _tool(mcp, "image_replace")(path=str(target), index=-1, src=str(new_png))


# ── PDF image extraction (best-effort) ──


def test_pdf_image_list_with_no_images(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    """A PDF with no embedded images returns count=0."""
    from reportlab.pdfgen import canvas

    docs_dir, _ = tmp_roots
    p = docs_dir / "noimg.pdf"
    c = canvas.Canvas(str(p))
    c.drawString(100, 750, "text only")
    c.save()
    mcp, _ctx = _setup(make_config())
    result = _tool(mcp, "image_list")(path=str(p))
    assert result["format"] == "pdf"
    assert result["count"] == 0


def test_pdf_image_list_corrupted_raises(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    from dokumen_pintar.errors import HandlerError

    docs_dir, _ = tmp_roots
    p = docs_dir / "broken.pdf"
    p.write_bytes(b"%PDF-not-real")
    mcp, _ctx = _setup(make_config())
    with pytest.raises(HandlerError, match="cannot read PDF"):
        _tool(mcp, "image_list")(path=str(p))


def test_pdf_image_extract_index_out_of_range(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    from reportlab.pdfgen import canvas

    docs_dir, _ = tmp_roots
    p = docs_dir / "noimg2.pdf"
    c = canvas.Canvas(str(p))
    c.drawString(100, 750, "x")
    c.save()
    out = docs_dir / "out.png"
    mcp, _ctx = _setup(make_config())
    with pytest.raises(ValidationError, match=r"(out of range|not found)"):
        _tool(mcp, "image_extract")(path=str(p), index=0, dst=str(out))


def test_pdf_extract_image_helper_not_found() -> None:
    """``_read_pdf_image_bytes`` raises when the requested index doesn't exist."""
    from reportlab.pdfgen import canvas
    import tempfile

    from dokumen_pintar.tools.images import _read_pdf_image_bytes

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
        c = canvas.Canvas(f.name)
        c.drawString(100, 750, "no images")
        c.save()
        target = Path(f.name)
    try:
        with pytest.raises(ValidationError, match="not found"):
            _read_pdf_image_bytes(target, 0)
    finally:
        target.unlink(missing_ok=True)


def test_replace_zip_member_keeps_archive_valid(tmp_path: Path) -> None:
    """``_replace_zip_member`` rebuilds a valid ZIP after swap."""
    from dokumen_pintar.tools.images import _replace_zip_member

    p = tmp_path / "test.zip"
    with zipfile.ZipFile(p, "w") as z:
        z.writestr("a.txt", "old-a")
        z.writestr("b.txt", "old-b")
    _replace_zip_member(p, "a.txt", b"new-a")
    with zipfile.ZipFile(p) as z:
        assert z.read("a.txt") == b"new-a"
        assert z.read("b.txt") == b"old-b"


def test_image_extract_all_pdf_with_no_images(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    """extract_all on a PDF without images writes nothing successfully."""
    from reportlab.pdfgen import canvas

    docs_dir, _ = tmp_roots
    p = docs_dir / "empty_pdf.pdf"
    c = canvas.Canvas(str(p))
    c.drawString(100, 750, "no images")
    c.save()
    out = docs_dir / "out_empty"
    mcp, _ctx = _setup(make_config())
    result = _tool(mcp, "image_extract_all")(path=str(p), dst_dir=str(out))
    assert result["count"] == 0
    assert out.exists() and out.is_dir()



# ── coverage gap fillers ──


def test_list_zip_skips_directory_entries(tmp_path: Path) -> None:
    """ZIP archive with a media-prefix directory entry returns empty list."""
    from dokumen_pintar.tools.images import _list_zip_images

    p = tmp_path / "dir_only.zip"
    with zipfile.ZipFile(p, "w") as z:
        # Directory entry under word/media/
        z.writestr("word/media/", b"")
    assert _list_zip_images(p, ("word/media/",)) == []


def test_list_zip_filters_non_image_extensions(tmp_path: Path) -> None:
    """Audio/video files under media/ are filtered out by extension."""
    from dokumen_pintar.tools.images import _list_zip_images

    p = tmp_path / "av.zip"
    with zipfile.ZipFile(p, "w") as z:
        z.writestr("word/media/audio.mp3", b"fake")
        z.writestr("word/media/img.png", _png_bytes())
    result = _list_zip_images(p, ("word/media/",))
    assert len(result) == 1
    assert result[0]["ext"] == ".png"


def test_pdf_image_list_encrypted_no_password(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    """Encrypted PDF with empty password decryption failure raises HandlerError."""
    from unittest.mock import MagicMock, PropertyMock, patch

    from dokumen_pintar.errors import HandlerError

    docs_dir, _ = tmp_roots
    p = docs_dir / "enc.pdf"
    # Make a stub PDF; we'll mock pypdf to claim it's encrypted.
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(str(p))
    c.drawString(100, 750, "x")
    c.save()

    fake_reader = MagicMock()
    type(fake_reader).is_encrypted = PropertyMock(return_value=True)
    fake_reader.decrypt.return_value = 0  # decrypt failed
    mcp, _ctx = _setup(make_config())
    with patch("pypdf.PdfReader", return_value=fake_reader):
        with pytest.raises(HandlerError, match="encrypted"):
            _tool(mcp, "image_list")(path=str(p))


def test_pdf_image_list_encrypted_decrypt_exception(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    """If pypdf.decrypt raises, we still surface HandlerError(encrypted)."""
    from unittest.mock import MagicMock, PropertyMock, patch

    from dokumen_pintar.errors import HandlerError

    docs_dir, _ = tmp_roots
    p = docs_dir / "encx.pdf"
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(str(p))
    c.drawString(100, 750, "y")
    c.save()
    fake_reader = MagicMock()
    type(fake_reader).is_encrypted = PropertyMock(return_value=True)
    fake_reader.decrypt.side_effect = Exception("crypto err")
    mcp, _ctx = _setup(make_config())
    with patch("pypdf.PdfReader", return_value=fake_reader):
        with pytest.raises(HandlerError, match="encrypted"):
            _tool(mcp, "image_list")(path=str(p))


def test_pdf_image_list_page_iteration_error_skipped(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    """A page whose .images iterator raises is silently skipped."""
    from unittest.mock import MagicMock, PropertyMock, patch

    docs_dir, _ = tmp_roots
    p = docs_dir / "skiperr.pdf"
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(str(p))
    c.drawString(100, 750, "z")
    c.save()
    bad_page = MagicMock()
    type(bad_page).images = PropertyMock(side_effect=Exception("page broken"))
    fake_reader = MagicMock()
    type(fake_reader).is_encrypted = PropertyMock(return_value=False)
    fake_reader.pages = [bad_page]
    mcp, _ctx = _setup(make_config())
    with patch("pypdf.PdfReader", return_value=fake_reader):
        result = _tool(mcp, "image_list")(path=str(p))
    assert result["count"] == 0


def test_read_pdf_image_bytes_skips_page_exception(tmp_path: Path) -> None:
    """``_read_pdf_image_bytes`` continues past pages whose .images raises."""
    from unittest.mock import MagicMock, PropertyMock, patch

    from dokumen_pintar.tools.images import _read_pdf_image_bytes

    p = tmp_path / "pgerr.pdf"
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(str(p))
    c.drawString(100, 750, "p")
    c.save()
    bad_page = MagicMock()
    type(bad_page).images = PropertyMock(side_effect=Exception("nope"))
    fake_reader = MagicMock()
    type(fake_reader).is_encrypted = PropertyMock(return_value=False)
    fake_reader.pages = [bad_page]
    with patch("pypdf.PdfReader", return_value=fake_reader):
        with pytest.raises(ValidationError, match="not found"):
            _read_pdf_image_bytes(p, 0)


def test_read_pdf_image_bytes_encrypted_decrypt_exception(tmp_path: Path) -> None:
    from unittest.mock import MagicMock, PropertyMock, patch

    from dokumen_pintar.errors import HandlerError
    from dokumen_pintar.tools.images import _read_pdf_image_bytes

    p = tmp_path / "encrypted.pdf"
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(str(p))
    c.drawString(100, 750, "e")
    c.save()
    fake_reader = MagicMock()
    type(fake_reader).is_encrypted = PropertyMock(return_value=True)
    fake_reader.decrypt.side_effect = Exception("nope")
    with patch("pypdf.PdfReader", return_value=fake_reader):
        with pytest.raises(HandlerError, match="encrypted"):
            _read_pdf_image_bytes(p, 0)


def test_pdf_image_extract_returns_match_via_helper(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    """Cover the path where _read_pdf_image_bytes finds a match (line 363->exit)."""
    from unittest.mock import MagicMock, PropertyMock, patch

    docs_dir, _ = tmp_roots
    p = docs_dir / "match.pdf"
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(str(p))
    c.drawString(100, 750, "m")
    c.save()
    out = docs_dir / "out.png"
    fake_img = MagicMock()
    fake_img.name = "im0.png"
    fake_img.data = _png_bytes()
    fake_page = MagicMock()
    type(fake_page).images = PropertyMock(return_value=[fake_img])
    fake_reader = MagicMock()
    type(fake_reader).is_encrypted = PropertyMock(return_value=False)
    fake_reader.pages = [fake_page]
    mcp, _ctx = _setup(make_config())
    with patch("pypdf.PdfReader", return_value=fake_reader):
        result = _tool(mcp, "image_extract")(path=str(p), index=0, dst=str(out))
    assert Path(result["dst"]).exists()
    assert Path(result["dst"]).stat().st_size > 0



def test_read_pdf_image_bytes_iterates_past_earlier_images(tmp_path: Path) -> None:
    """When index > 0, the helper must increment ``seen`` for earlier images."""
    from unittest.mock import MagicMock, PropertyMock, patch

    from dokumen_pintar.tools.images import _read_pdf_image_bytes

    p = tmp_path / "multi.pdf"
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(str(p))
    c.drawString(100, 750, "x")
    c.save()

    img0 = MagicMock(name="i0")
    img0.name = "img0.png"
    img0.data = _png_bytes("red")
    img1 = MagicMock(name="i1")
    img1.name = "img1.png"
    img1.data = _png_bytes("blue")
    fake_page = MagicMock()
    type(fake_page).images = PropertyMock(return_value=[img0, img1])
    fake_reader = MagicMock()
    type(fake_reader).is_encrypted = PropertyMock(return_value=False)
    fake_reader.pages = [fake_page]

    with patch("pypdf.PdfReader", return_value=fake_reader):
        blob = _read_pdf_image_bytes(p, 1)
    assert blob == img1.data


def test_pdf_image_list_encrypted_decrypt_succeeds(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    """Encrypted PDF whose empty-password decrypt returns truthy proceeds."""
    from unittest.mock import MagicMock, PropertyMock, patch

    docs_dir, _ = tmp_roots
    p = docs_dir / "encok.pdf"
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(str(p))
    c.drawString(100, 750, "x")
    c.save()
    fake_reader = MagicMock()
    type(fake_reader).is_encrypted = PropertyMock(return_value=True)
    fake_reader.decrypt.return_value = 1
    fake_reader.pages = []
    mcp, _ctx = _setup(make_config())
    with patch("pypdf.PdfReader", return_value=fake_reader):
        result = _tool(mcp, "image_list")(path=str(p))
    assert result["count"] == 0



def test_replace_zip_member_finally_skips_when_tmp_consumed(tmp_path: Path) -> None:
    """When the rebuild path completes and tmp is removed mid-finally,
    the cleanup branch must take the false-skip exit (line 363->exit)."""
    from unittest.mock import patch

    from dokumen_pintar.tools.images import _replace_zip_member

    p = tmp_path / "tc.zip"
    with zipfile.ZipFile(p, "w") as z:
        z.writestr("a.txt", "v1")

    real_exists = Path.exists
    call_count = {"n": 0}

    def patched(self, *a, **kw):  # type: ignore[no-untyped-def]
        # The finally block does Path(...).exists() on the .tmp file - first
        # call returns False so the unlink branch is skipped entirely.
        if self.suffix == ".tmp":
            call_count["n"] += 1
            return False
        return real_exists(self, *a, **kw)

    with patch.object(Path, "exists", patched):
        _replace_zip_member(p, "a.txt", b"v2")
    assert call_count["n"] >= 1
    with zipfile.ZipFile(p) as z:
        assert z.read("a.txt") == b"v2"
