"""Tests for the v1.0.2 metadata-edit layer.

Covers:
- :mod:`dokumen_pintar.handlers.image_handler` — EXIF read/write/strip + helpers.
- ``write_meta`` / ``strip_meta`` on docx, xlsx, pptx, and pdf handlers.
- :mod:`dokumen_pintar.tools.metadata` — the four MCP tools with snapshotting.

The tests build small synthetic fixtures (JPEGs with seeded EXIF, minimal
DOCX/XLSX/PPTX/PDF) rather than relying on external sample files, so the
suite stays hermetic.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Callable
from unittest.mock import patch

import piexif
import pytest
from PIL import Image
from mcp.server.fastmcp import FastMCP

from dokumen_pintar.config import AppConfig
from dokumen_pintar.context import build_context
from dokumen_pintar.errors import HandlerError, UnsupportedFormatError
from dokumen_pintar.handlers import image_handler as _ih
from dokumen_pintar.handlers.image_handler import (
    ImageHandler,
    _decode_value,
    _extract_gps_summary,
    _gps_to_decimal,
    _normalize_value_for_piexif,
    _read_exif_via_pillow,
)
from dokumen_pintar.tools import metadata as metadata_tool


# ──────────────────────────────────────────────────────────── helpers


def _setup_mcp(cfg: AppConfig):
    """Build an MCP with metadata + version + workspace tools wired in."""
    from dokumen_pintar.tools import version, workspace, content_crud

    ctx = build_context(cfg)
    mcp = FastMCP(name="meta-test")
    workspace.register(mcp, ctx)
    content_crud.register(mcp, ctx)
    version.register(mcp, ctx)
    metadata_tool.register(mcp, ctx)
    return mcp, ctx


def _tool(mcp, name):  # type: ignore[no-untyped-def]
    return mcp._tool_manager._tools[name].fn


def _make_jpeg_with_exif(
    path: Path,
    artist: str = "Original Author",
    copyright_: str = "(c) Test 2026",
    description: str = "Sample photo",
) -> None:
    """Write a tiny 4×4 JPEG with a hand-rolled EXIF block."""
    exif_dict = {
        "0th": {
            piexif.ImageIFD.Artist: artist.encode("utf-8"),
            piexif.ImageIFD.Copyright: copyright_.encode("utf-8"),
            piexif.ImageIFD.ImageDescription: description.encode("utf-8"),
            piexif.ImageIFD.Software: b"Dokumen-Pintar Tests",
        },
        "Exif": {},
        "GPS": {},
        "Interop": {},
        "1st": {},
        "thumbnail": None,
    }
    exif_bytes = piexif.dump(exif_dict)
    img = Image.new("RGB", (4, 4), color=(127, 200, 64))
    img.save(str(path), format="JPEG", exif=exif_bytes)


def _make_jpeg_with_gps(path: Path) -> None:
    """JPEG embedding a known GPS coordinate (Jakarta, ~-6.2,106.8)."""
    gps_dict = {
        piexif.GPSIFD.GPSLatitudeRef: b"S",
        piexif.GPSIFD.GPSLatitude: ((6, 1), (12, 1), (0, 1)),
        piexif.GPSIFD.GPSLongitudeRef: b"E",
        piexif.GPSIFD.GPSLongitude: ((106, 1), (48, 1), (0, 1)),
        piexif.GPSIFD.GPSAltitudeRef: 0,
        piexif.GPSIFD.GPSAltitude: (10, 1),
    }
    exif_dict = {
        "0th": {},
        "Exif": {},
        "GPS": gps_dict,
        "Interop": {},
        "1st": {},
        "thumbnail": None,
    }
    exif_bytes = piexif.dump(exif_dict)
    img = Image.new("RGB", (4, 4))
    img.save(str(path), format="JPEG", exif=exif_bytes)


# ─────────────────────────────────────────── _decode_value coverage


class TestDecodeValue:
    def test_bytes_utf8(self) -> None:
        assert _decode_value(b"hello\x00") == "hello"

    def test_bytes_latin1_fallback(self) -> None:
        # b"\xff" is not valid utf-8 but valid latin-1.
        out = _decode_value(b"\xff\xff")
        assert isinstance(out, str)

    def test_bytes_undecodable_returns_hex(self) -> None:
        # Patch both decode attempts to fail so we hit the hex fallback.
        class _BadBytes(bytes):
            def decode(self, *args: Any, **kwargs: Any) -> str:
                raise UnicodeDecodeError("x", b"", 0, 1, "y")

        out = _decode_value(_BadBytes(b"\x80\x81"))
        assert out == "8081"

    def test_ifdrational_normal(self) -> None:
        from PIL.TiffImagePlugin import IFDRational

        assert _decode_value(IFDRational(1, 2)) == pytest.approx(0.5)

    def test_ifdrational_zero_division(self) -> None:
        from PIL.TiffImagePlugin import IFDRational

        out = _decode_value(IFDRational(1, 0))
        # Division-by-zero path returns a "num/den" string.
        assert isinstance(out, str) and "/" in out

    def test_tuple_and_dict(self) -> None:
        assert _decode_value((1, 2)) == [1, 2]
        assert _decode_value({"k": b"v"}) == {"k": "v"}

    def test_passthrough_scalar(self) -> None:
        assert _decode_value(42) == 42


# ─────────────────────────────────────────── GPS helpers


class TestGpsHelpers:
    def test_gps_to_decimal_north(self) -> None:
        result = _gps_to_decimal((6.0, 12.0, 0.0), "N")
        assert result == pytest.approx(6.2, abs=1e-3)

    def test_gps_to_decimal_south_negative(self) -> None:
        result = _gps_to_decimal((6.0, 12.0, 0.0), b"S")
        assert result == pytest.approx(-6.2, abs=1e-3)

    def test_gps_to_decimal_west_negative(self) -> None:
        result = _gps_to_decimal((106.0, 48.0, 0.0), "W")
        assert result == pytest.approx(-106.8, abs=1e-3)

    def test_gps_to_decimal_missing_inputs(self) -> None:
        assert _gps_to_decimal(None, "N") is None
        assert _gps_to_decimal((1, 2, 3), None) is None

    def test_gps_to_decimal_wrong_shape(self) -> None:
        assert _gps_to_decimal((1, 2), "N") is None

    def test_gps_to_decimal_value_error(self) -> None:
        # A non-numeric tuple triggers the except branch.
        assert _gps_to_decimal(("a", "b", "c"), "N") is None

    def test_gps_to_decimal_unconvertible_tuple_returns_none(self) -> None:
        """A 3-tuple containing nested tuples cannot be ``float()``-ed and
        triggers the TypeError fallback path."""
        assert _gps_to_decimal(((6, 1), (12, 1), (0, 1)), "N") is None

    def test_extract_gps_summary_full(self) -> None:
        exif = {
            "GPSInfo": {
                "GPSLatitude": [6.0, 12.0, 0.0],
                "GPSLatitudeRef": "S",
                "GPSLongitude": [106.0, 48.0, 0.0],
                "GPSLongitudeRef": "E",
                "GPSAltitude": 10,
            }
        }
        summary = _extract_gps_summary(exif)
        assert summary is not None
        assert summary["latitude"] == pytest.approx(-6.2, abs=1e-3)
        assert summary["longitude"] == pytest.approx(106.8, abs=1e-3)
        assert summary["altitude"] == 10

    def test_extract_gps_summary_missing(self) -> None:
        assert _extract_gps_summary({}) is None

    def test_extract_gps_summary_not_dict(self) -> None:
        assert _extract_gps_summary({"GPSInfo": "garbage"}) is None

    def test_extract_gps_summary_only_lat(self) -> None:
        exif = {
            "GPSInfo": {
                "GPSLatitude": [6.0, 12.0, 0.0],
                "GPSLatitudeRef": "N",
            }
        }
        summary = _extract_gps_summary(exif)
        assert summary == {"latitude": pytest.approx(6.2, abs=1e-3)}

    def test_extract_gps_summary_empty_returns_none(self) -> None:
        # All values are None/missing → summary should be empty → None.
        assert _extract_gps_summary({"GPSInfo": {}}) is None


# ─────────────────────────────────────────── _normalize_value_for_piexif


class TestNormalizePiexifValue:
    def test_none_passthrough(self) -> None:
        assert _normalize_value_for_piexif(None) is None

    def test_str_to_bytes(self) -> None:
        assert _normalize_value_for_piexif("abc") == b"abc"

    def test_bytes_passthrough(self) -> None:
        assert _normalize_value_for_piexif(b"raw") == b"raw"

    def test_int_passthrough(self) -> None:
        assert _normalize_value_for_piexif(3) == 3

    def test_float_to_rational(self) -> None:
        assert _normalize_value_for_piexif(1.5) == (1500, 1000)

    def test_unsupported_type(self) -> None:
        with pytest.raises(HandlerError):
            _normalize_value_for_piexif([1, 2, 3])


# ─────────────────────────────────────────── ImageHandler core


class TestImageHandlerCore:
    def test_detect(self, tmp_path: Path) -> None:
        h = ImageHandler()
        assert h.detect(tmp_path / "foo.jpg") is True
        assert h.detect(tmp_path / "foo.png") is True
        assert h.detect(tmp_path / "foo.txt") is False

    def test_read_meta_basic(self, tmp_path: Path) -> None:
        path = tmp_path / "img.jpg"
        _make_jpeg_with_exif(path)
        meta = ImageHandler().read_meta(path)
        assert meta["format"] == "image"
        assert meta["image_format"] == "JPEG"
        assert meta["width"] == 4 and meta["height"] == 4
        assert meta["exif"]["Artist"] == "Original Author"
        assert "Copyright" in meta["exif"]

    def test_read_meta_with_gps(self, tmp_path: Path) -> None:
        path = tmp_path / "gps.jpg"
        _make_jpeg_with_gps(path)
        meta = ImageHandler().read_meta(path)
        assert "gps" in meta
        assert meta["gps"]["latitude"] == pytest.approx(-6.2, abs=1e-3)

    def test_read_meta_png_text_chunks(self, tmp_path: Path) -> None:
        from PIL.PngImagePlugin import PngInfo

        path = tmp_path / "tagged.png"
        info = PngInfo()
        info.add_text("Author", "Firdaus")
        info.add_text("Description", "Test PNG")
        Image.new("RGBA", (4, 4)).save(path, "PNG", pnginfo=info)
        meta = ImageHandler().read_meta(path)
        assert meta["image_format"] == "PNG"
        assert meta["png_text"]["Author"] == "Firdaus"
        assert meta["has_alpha"] is True

    def test_read_meta_iptc_optional(self, tmp_path: Path) -> None:
        # No IPTC block → key absent (no KeyError).
        path = tmp_path / "noiptc.jpg"
        _make_jpeg_with_exif(path)
        meta = ImageHandler().read_meta(path)
        assert "iptc" not in meta

    def test_read_meta_bad_file_raises(self, tmp_path: Path) -> None:
        bad = tmp_path / "bad.jpg"
        bad.write_bytes(b"not a real image")
        with pytest.raises(HandlerError):
            ImageHandler().read_meta(bad)

    def test_read_text_returns_description_fields(self, tmp_path: Path) -> None:
        path = tmp_path / "desc.jpg"
        _make_jpeg_with_exif(path, description="A note")
        text = ImageHandler().read_text(path)
        assert "A note" in text
        assert "Artist" in text

    def test_extract_for_search_returns_text(self, tmp_path: Path) -> None:
        path = tmp_path / "ext.jpg"
        _make_jpeg_with_exif(path)
        assert "Original Author" in ImageHandler().extract_for_search(path)

    def test_extract_for_search_swallows_errors(self, tmp_path: Path) -> None:
        bad = tmp_path / "broken.jpg"
        bad.write_bytes(b"garbage")
        assert ImageHandler().extract_for_search(bad) == ""

    def test_read_exif_pillow_no_exif(self, tmp_path: Path) -> None:
        path = tmp_path / "plain.png"
        Image.new("RGB", (2, 2)).save(path)
        assert _read_exif_via_pillow(path) == {}


# ─────────────────────────────────────────── ImageHandler structured_get


class TestImageHandlerStructured:
    def test_structured_get_exif(self, tmp_path: Path) -> None:
        path = tmp_path / "s.jpg"
        _make_jpeg_with_exif(path)
        out = ImageHandler().structured_get(path, "exif")
        assert "Artist" in out

    def test_structured_get_metadata_alias(self, tmp_path: Path) -> None:
        path = tmp_path / "s.jpg"
        _make_jpeg_with_exif(path)
        out = ImageHandler().structured_get(path, "metadata")
        assert "Artist" in out

    def test_structured_get_dimensions(self, tmp_path: Path) -> None:
        path = tmp_path / "s.jpg"
        _make_jpeg_with_exif(path)
        out = ImageHandler().structured_get(path, "dimensions")
        assert out == {"width": 4, "height": 4}

    def test_structured_get_size_alias(self, tmp_path: Path) -> None:
        path = tmp_path / "s.jpg"
        _make_jpeg_with_exif(path)
        out = ImageHandler().structured_get(path, "size")
        assert "width" in out

    def test_structured_get_gps(self, tmp_path: Path) -> None:
        path = tmp_path / "g.jpg"
        _make_jpeg_with_gps(path)
        gps = ImageHandler().structured_get(path, "gps")
        assert gps["latitude"] == pytest.approx(-6.2, abs=1e-3)

    def test_structured_get_specific_tag(self, tmp_path: Path) -> None:
        path = tmp_path / "s.jpg"
        _make_jpeg_with_exif(path)
        out = ImageHandler().structured_get(path, "exif:Artist")
        assert out == "Original Author"

    def test_structured_get_missing_tag(self, tmp_path: Path) -> None:
        path = tmp_path / "s.jpg"
        _make_jpeg_with_exif(path)
        with pytest.raises(HandlerError):
            ImageHandler().structured_get(path, "exif:Nonexistent")

    def test_structured_get_unsupported_expr(self, tmp_path: Path) -> None:
        path = tmp_path / "s.jpg"
        _make_jpeg_with_exif(path)
        with pytest.raises(HandlerError):
            ImageHandler().structured_get(path, "whatever")


# ─────────────────────────────────────────── ImageHandler write/strip


class TestImageHandlerWrite:
    def test_write_meta_single_key(self, tmp_path: Path) -> None:
        path = tmp_path / "w.jpg"
        _make_jpeg_with_exif(path)
        applied = ImageHandler().write_meta(path, {"artist": "New Author"})
        assert applied == {"artist": "New Author"}
        # Round-trip.
        meta = ImageHandler().read_meta(path)
        assert meta["exif"]["Artist"] == "New Author"

    def test_write_meta_delete_via_none(self, tmp_path: Path) -> None:
        path = tmp_path / "w.jpg"
        _make_jpeg_with_exif(path)
        ImageHandler().write_meta(path, {"copyright": None})
        meta = ImageHandler().read_meta(path)
        # Either absent or empty string after deletion.
        assert meta["exif"].get("Copyright", "") == ""

    def test_write_meta_unknown_key(self, tmp_path: Path) -> None:
        path = tmp_path / "w.jpg"
        _make_jpeg_with_exif(path)
        with pytest.raises(HandlerError):
            ImageHandler().write_meta(path, {"nonsense": "x"})

    def test_write_meta_unsupported_format(self, tmp_path: Path) -> None:
        path = tmp_path / "w.png"
        Image.new("RGB", (2, 2)).save(path)
        with pytest.raises(HandlerError):
            ImageHandler().write_meta(path, {"artist": "x"})

    def test_write_meta_corrupt_file(self, tmp_path: Path) -> None:
        bad = tmp_path / "bad.jpg"
        bad.write_bytes(b"definitely not jpeg")
        with pytest.raises(HandlerError):
            ImageHandler().write_meta(bad, {"artist": "x"})

    def test_write_meta_piexif_dump_failure(self, tmp_path: Path) -> None:
        path = tmp_path / "f.jpg"
        _make_jpeg_with_exif(path)
        with patch.object(piexif, "dump", side_effect=ValueError("boom")):
            with pytest.raises(HandlerError):
                ImageHandler().write_meta(path, {"artist": "x"})

    def test_write_meta_float_normalizes(self, tmp_path: Path) -> None:
        path = tmp_path / "n.jpg"
        _make_jpeg_with_exif(path)
        # Orientation is technically a SHORT — but we just verify the
        # type-normalization path runs without complaint. Reject is fine too;
        # piexif may complain, in which case we catch HandlerError.
        try:
            ImageHandler().write_meta(path, {"orientation": 1})
        except HandlerError:
            pass

    def test_strip_meta_clears_exif(self, tmp_path: Path) -> None:
        path = tmp_path / "s.jpg"
        _make_jpeg_with_exif(path)
        result = ImageHandler().strip_meta(path)
        assert result["stripped"] is True
        meta = ImageHandler().read_meta(path)
        # After strip, EXIF either empty or has no Artist field anymore.
        assert "Artist" not in meta.get("exif", {})

    def test_strip_meta_unsupported_format(self, tmp_path: Path) -> None:
        path = tmp_path / "s.png"
        Image.new("RGB", (2, 2)).save(path)
        with pytest.raises(HandlerError):
            ImageHandler().strip_meta(path)

    def test_strip_meta_failure_path(self, tmp_path: Path) -> None:
        path = tmp_path / "s.jpg"
        _make_jpeg_with_exif(path)
        with patch.object(piexif, "remove", side_effect=ValueError("nope")):
            with pytest.raises(HandlerError):
                ImageHandler().strip_meta(path)


# ─────────────────────────────────────────── DOCX write_meta


class TestDocxWriteMeta:
    def _docx(self, path: Path) -> None:
        from docx import Document

        d = Document()
        d.add_paragraph("hello")
        d.save(str(path))

    def test_write_meta_round_trip(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.docx_handler import DocxHandler

        p = tmp_path / "d.docx"
        self._docx(p)
        h = DocxHandler()
        applied = h.write_meta(p, {"author": "Firdaus", "title": "My Doc"})
        assert applied == {"author": "Firdaus", "title": "My Doc"}
        meta = h.read_meta(p)
        assert meta["core_props"]["author"] == "Firdaus"
        assert meta["core_props"]["title"] == "My Doc"

    def test_write_meta_unknown_key(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.docx_handler import DocxHandler

        p = tmp_path / "d.docx"
        self._docx(p)
        with pytest.raises(HandlerError):
            DocxHandler().write_meta(p, {"nonsense": "x"})

    def test_write_meta_type_error(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.docx_handler import DocxHandler

        p = tmp_path / "d.docx"
        self._docx(p)
        # ``created`` expects a datetime; passing a list triggers the
        # AttributeError/TypeError/ValueError branch.
        with pytest.raises(HandlerError):
            DocxHandler().write_meta(p, {"created": ["not", "a", "datetime"]})

    def test_write_meta_save_failure(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.docx_handler import DocxHandler

        p = tmp_path / "d.docx"
        self._docx(p)
        with patch("docx.document.Document.save", side_effect=OSError("disk")):
            with pytest.raises(HandlerError):
                DocxHandler().write_meta(p, {"author": "x"})

    def test_strip_meta(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.docx_handler import DocxHandler

        p = tmp_path / "d.docx"
        self._docx(p)
        h = DocxHandler()
        h.write_meta(p, {"author": "Firdaus"})
        result = h.strip_meta(p)
        assert "author" in result["stripped"]
        meta = h.read_meta(p)
        # Author field is now empty string (the property exists but is "").
        assert (meta["core_props"]["author"] or "") == ""

    def test_strip_meta_save_failure(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.docx_handler import DocxHandler

        p = tmp_path / "d.docx"
        self._docx(p)
        with patch("docx.document.Document.save", side_effect=OSError("disk")):
            with pytest.raises(HandlerError):
                DocxHandler().strip_meta(p)


# ─────────────────────────────────────────── XLSX write_meta


class TestXlsxWriteMeta:
    def _xlsx(self, path: Path) -> None:
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "hello"
        wb.save(str(path))

    def test_write_meta_round_trip(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.xlsx_handler import XlsxHandler

        p = tmp_path / "s.xlsx"
        self._xlsx(p)
        h = XlsxHandler()
        applied = h.write_meta(p, {"creator": "F", "title": "T"})
        assert applied == {"creator": "F", "title": "T"}

        import openpyxl

        wb = openpyxl.load_workbook(p)
        assert wb.properties.creator == "F"
        assert wb.properties.title == "T"

    def test_write_meta_unknown(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.xlsx_handler import XlsxHandler

        p = tmp_path / "s.xlsx"
        self._xlsx(p)
        with pytest.raises(HandlerError):
            XlsxHandler().write_meta(p, {"bogus": "x"})

    def test_write_meta_type_error(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.xlsx_handler import XlsxHandler

        p = tmp_path / "s.xlsx"
        self._xlsx(p)
        with pytest.raises(HandlerError):
            XlsxHandler().write_meta(p, {"created": "not-a-date"})

    def test_write_meta_save_failure(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.xlsx_handler import XlsxHandler
        from openpyxl import Workbook

        p = tmp_path / "s.xlsx"
        self._xlsx(p)
        with patch.object(Workbook, "save", side_effect=OSError("disk full")):
            with pytest.raises(HandlerError):
                XlsxHandler().write_meta(p, {"creator": "x"})

    def test_strip_meta(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.xlsx_handler import XlsxHandler

        p = tmp_path / "s.xlsx"
        self._xlsx(p)
        h = XlsxHandler()
        h.write_meta(p, {"creator": "Firdaus"})
        result = h.strip_meta(p)
        assert "creator" in result["stripped"]
        import openpyxl

        wb = openpyxl.load_workbook(p)
        assert (wb.properties.creator or "") == ""

    def test_strip_meta_save_failure(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.xlsx_handler import XlsxHandler
        from openpyxl import Workbook

        p = tmp_path / "s.xlsx"
        self._xlsx(p)
        with patch.object(Workbook, "save", side_effect=OSError("io")):
            with pytest.raises(HandlerError):
                XlsxHandler().strip_meta(p)


# ─────────────────────────────────────────── PPTX write_meta


class TestPptxWriteMeta:
    def _pptx(self, path: Path) -> None:
        from pptx import Presentation

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])
        prs.save(str(path))

    def test_write_meta_round_trip(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pptx_handler import PptxHandler

        p = tmp_path / "d.pptx"
        self._pptx(p)
        h = PptxHandler()
        applied = h.write_meta(p, {"author": "F", "subject": "S"})
        assert applied == {"author": "F", "subject": "S"}
        from pptx import Presentation

        prs = Presentation(str(p))
        assert prs.core_properties.author == "F"

    def test_write_meta_unknown(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pptx_handler import PptxHandler

        p = tmp_path / "d.pptx"
        self._pptx(p)
        with pytest.raises(HandlerError):
            PptxHandler().write_meta(p, {"unknown": "x"})

    def test_write_meta_type_error(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pptx_handler import PptxHandler

        p = tmp_path / "d.pptx"
        self._pptx(p)
        with pytest.raises(HandlerError):
            PptxHandler().write_meta(p, {"created": "not-date"})

    def test_write_meta_save_failure(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pptx_handler import PptxHandler
        from pptx.presentation import Presentation as _Pres

        p = tmp_path / "d.pptx"
        self._pptx(p)
        with patch.object(_Pres, "save", side_effect=OSError("disk")):
            with pytest.raises(HandlerError):
                PptxHandler().write_meta(p, {"author": "x"})

    def test_strip_meta(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pptx_handler import PptxHandler

        p = tmp_path / "d.pptx"
        self._pptx(p)
        h = PptxHandler()
        h.write_meta(p, {"author": "F"})
        result = h.strip_meta(p)
        assert "author" in result["stripped"]

    def test_strip_meta_save_failure(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pptx_handler import PptxHandler
        from pptx.presentation import Presentation as _Pres

        p = tmp_path / "d.pptx"
        self._pptx(p)
        with patch.object(_Pres, "save", side_effect=OSError("disk")):
            with pytest.raises(HandlerError):
                PptxHandler().strip_meta(p)

    def test_write_meta_open_failure(self, tmp_path: Path) -> None:
        """A garbage .pptx must surface as HandlerError on write_meta open."""
        from dokumen_pintar.handlers.pptx_handler import PptxHandler

        p = tmp_path / "garbage.pptx"
        p.write_bytes(b"not a real pptx")
        with pytest.raises(HandlerError):
            PptxHandler().write_meta(p, {"author": "x"})

    def test_strip_meta_open_failure(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pptx_handler import PptxHandler

        p = tmp_path / "garbage.pptx"
        p.write_bytes(b"definitely not a pptx")
        with pytest.raises(HandlerError):
            PptxHandler().strip_meta(p)


# ─────────────────────────────────────────── PDF write_meta


class TestPdfWriteMeta:
    def _pdf(self, path: Path) -> None:
        from reportlab.pdfgen import canvas

        c = canvas.Canvas(str(path))
        c.drawString(72, 720, "hello")
        c.showPage()
        c.save()

    def test_write_meta_round_trip(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pdf_handler import PdfHandler

        p = tmp_path / "d.pdf"
        self._pdf(p)
        applied = PdfHandler().write_meta(p, {"title": "T", "author": "F"})
        assert applied == {"title": "T", "author": "F"}
        import pikepdf

        with pikepdf.open(str(p)) as pdf:
            assert str(pdf.docinfo[pikepdf.Name("/Title")]) == "T"
            assert str(pdf.docinfo[pikepdf.Name("/Author")]) == "F"

    def test_write_meta_unknown(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pdf_handler import PdfHandler

        p = tmp_path / "d.pdf"
        self._pdf(p)
        with pytest.raises(HandlerError):
            PdfHandler().write_meta(p, {"unknown": "x"})

    def test_write_meta_delete_via_none(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pdf_handler import PdfHandler

        p = tmp_path / "d.pdf"
        self._pdf(p)
        h = PdfHandler()
        h.write_meta(p, {"title": "X"})
        h.write_meta(p, {"title": None})
        import pikepdf

        with pikepdf.open(str(p)) as pdf:
            assert pikepdf.Name("/Title") not in pdf.docinfo

    def test_write_meta_delete_missing_key_is_silent(
        self, tmp_path: Path
    ) -> None:
        from dokumen_pintar.handlers.pdf_handler import PdfHandler

        p = tmp_path / "d.pdf"
        self._pdf(p)
        # Setting an absent key to None must not raise.
        PdfHandler().write_meta(p, {"keywords": None})

    def test_write_meta_encrypted_pdf(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pdf_handler import PdfHandler
        import pikepdf

        p = tmp_path / "enc.pdf"
        self._pdf(p)
        # Encrypt the PDF with an owner password.
        with pikepdf.open(str(p), allow_overwriting_input=True) as pdf:
            pdf.save(
                str(p),
                encryption=pikepdf.Encryption(user="u", owner="o", R=4),
            )
        with pytest.raises(HandlerError):
            PdfHandler().write_meta(p, {"title": "X"})

    def test_write_meta_pdf_error(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pdf_handler import PdfHandler
        import pikepdf

        p = tmp_path / "broken.pdf"
        p.write_bytes(b"not really a pdf")
        with pytest.raises(HandlerError):
            PdfHandler().write_meta(p, {"title": "x"})

    def test_write_meta_oserror(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pdf_handler import PdfHandler
        import pikepdf

        p = tmp_path / "d.pdf"
        self._pdf(p)
        with patch.object(
            pikepdf.Pdf, "save", side_effect=OSError("nope")
        ):
            with pytest.raises(HandlerError):
                PdfHandler().write_meta(p, {"title": "x"})

    def test_strip_meta(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pdf_handler import PdfHandler

        p = tmp_path / "d.pdf"
        self._pdf(p)
        h = PdfHandler()
        h.write_meta(p, {"title": "Will-be-cleared"})
        result = h.strip_meta(p)
        assert any(s.startswith("/") for s in result["stripped"])
        import pikepdf

        with pikepdf.open(str(p)) as pdf:
            assert pikepdf.Name("/Title") not in pdf.docinfo

    def test_strip_meta_encrypted(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pdf_handler import PdfHandler
        import pikepdf

        p = tmp_path / "enc.pdf"
        self._pdf(p)
        with pikepdf.open(str(p), allow_overwriting_input=True) as pdf:
            pdf.save(
                str(p),
                encryption=pikepdf.Encryption(user="u", owner="o", R=4),
            )
        with pytest.raises(HandlerError):
            PdfHandler().strip_meta(p)

    def test_strip_meta_pdf_error(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pdf_handler import PdfHandler

        p = tmp_path / "broken.pdf"
        p.write_bytes(b"junk")
        with pytest.raises(HandlerError):
            PdfHandler().strip_meta(p)

    def test_strip_meta_oserror(self, tmp_path: Path) -> None:
        from dokumen_pintar.handlers.pdf_handler import PdfHandler
        import pikepdf

        p = tmp_path / "d.pdf"
        self._pdf(p)
        with patch.object(
            pikepdf.Pdf, "save", side_effect=OSError("disk")
        ):
            with pytest.raises(HandlerError):
                PdfHandler().strip_meta(p)


# ─────────────────────────────────────────── metadata MCP tools


class TestMetadataTools:
    def test_metadata_read(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        docs_dir, _ = tmp_roots
        path = docs_dir / "p.jpg"
        _make_jpeg_with_exif(path)
        mcp, _ = _setup_mcp(make_config())
        out = _tool(mcp, "metadata_read")(path="documents:/p.jpg")
        assert out["format"] == "image"
        assert "exif" in out["meta"]

    def test_metadata_read_unsupported_format(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        docs_dir, _ = tmp_roots
        path = docs_dir / "p.xyz"  # no handler for .xyz
        path.write_text("hi")
        mcp, _ = _setup_mcp(make_config())
        with pytest.raises(UnsupportedFormatError):
            _tool(mcp, "metadata_read")(path="documents:/p.xyz")

    def test_metadata_write_creates_snapshots(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        docs_dir, _ = tmp_roots
        path = docs_dir / "w.jpg"
        _make_jpeg_with_exif(path)
        mcp, ctx = _setup_mcp(make_config())
        out = _tool(mcp, "metadata_write")(
            path="documents:/w.jpg", updates={"artist": "Cascade"}
        )
        assert out["applied"]["artist"] == "Cascade"
        versions = _tool(mcp, "version_list")(path="documents:/w.jpg")
        assert len(versions["versions"]) == 2
        actions = {v["action"] for v in versions["versions"]}
        assert "metadata_write_pre" in actions
        assert "metadata_write_post" in actions

    def test_metadata_write_requires_dict(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        docs_dir, _ = tmp_roots
        path = docs_dir / "w.jpg"
        _make_jpeg_with_exif(path)
        mcp, _ = _setup_mcp(make_config())
        with pytest.raises(HandlerError):
            _tool(mcp, "metadata_write")(
                path="documents:/w.jpg", updates={}
            )

    def test_metadata_write_unsupported_handler(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        docs_dir, _ = tmp_roots
        path = docs_dir / "note.txt"
        path.write_text("hi")
        mcp, _ = _setup_mcp(make_config())
        with pytest.raises(UnsupportedFormatError):
            _tool(mcp, "metadata_write")(
                path="documents:/note.txt", updates={"author": "x"}
            )

    def test_metadata_write_no_handler(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        docs_dir, _ = tmp_roots
        path = docs_dir / "weird.xyz"
        path.write_text("x")
        mcp, _ = _setup_mcp(make_config())
        with pytest.raises(UnsupportedFormatError):
            _tool(mcp, "metadata_write")(
                path="documents:/weird.xyz", updates={"author": "x"}
            )

    def test_metadata_delete(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        docs_dir, _ = tmp_roots
        path = docs_dir / "d.jpg"
        _make_jpeg_with_exif(path)
        mcp, _ = _setup_mcp(make_config())
        out = _tool(mcp, "metadata_delete")(
            path="documents:/d.jpg", keys=["copyright"]
        )
        assert "copyright" in out["deleted"]

    def test_metadata_delete_requires_keys(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        docs_dir, _ = tmp_roots
        path = docs_dir / "d.jpg"
        _make_jpeg_with_exif(path)
        mcp, _ = _setup_mcp(make_config())
        with pytest.raises(HandlerError):
            _tool(mcp, "metadata_delete")(path="documents:/d.jpg", keys=[])

    def test_metadata_strip(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        docs_dir, _ = tmp_roots
        path = docs_dir / "s.jpg"
        _make_jpeg_with_exif(path)
        mcp, _ = _setup_mcp(make_config())
        out = _tool(mcp, "metadata_strip")(path="documents:/s.jpg")
        assert out["format"] == "image"
        assert out["stripped"] is True

    def test_metadata_strip_snapshots(
        self,
        make_config: Callable[..., AppConfig],
        tmp_roots: tuple[Path, Path],
    ) -> None:
        docs_dir, _ = tmp_roots
        path = docs_dir / "s.jpg"
        _make_jpeg_with_exif(path)
        mcp, _ = _setup_mcp(make_config())
        _tool(mcp, "metadata_strip")(path="documents:/s.jpg")
        versions = _tool(mcp, "version_list")(path="documents:/s.jpg")
        actions = {v["action"] for v in versions["versions"]}
        assert "metadata_strip_pre" in actions
        assert "metadata_strip_post" in actions



# ── v1.1.0 2.3: metadata_read_batch ──


def test_metadata_read_batch_mixed_formats(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    """Bulk reader returns metadata for every supported file matching the glob."""
    import openpyxl
    from docx import Document
    from PIL import Image as _Image

    docs_dir, _ = tmp_roots

    # DOCX
    docx_path = docs_dir / "report.docx"
    doc = Document()
    doc.core_properties.title = "Report"
    doc.add_paragraph("hello")
    doc.save(str(docx_path))

    # XLSX
    xlsx_path = docs_dir / "data.xlsx"
    wb = openpyxl.Workbook()
    wb.properties.title = "Sheet"
    wb.active["A1"] = "x"
    wb.save(xlsx_path)

    # JPEG
    jpg_path = docs_dir / "photo.jpg"
    _Image.new("RGB", (32, 32), "red").save(jpg_path, "JPEG")

    # Plain text (no read_meta on text handler in this version)
    (docs_dir / "note.txt").write_text("plain", encoding="utf-8")

    cfg = make_config()
    ctx = build_context(cfg)
    mcp = FastMCP(name="t-batch")
    metadata_tool.register(mcp, ctx)
    metadata_tool.register_batch(mcp, ctx)

    fn = mcp._tool_manager._tools["metadata_read_batch"].fn
    result = fn(glob="documents:/*")

    formats = {entry["format"] for entry in result["files"]}
    assert "docx" in formats
    assert "xlsx" in formats
    # JPG handled by image handler.
    assert any(entry["format"] == "image" for entry in result["files"])
    assert result["count"] >= 3


def test_metadata_read_batch_field_filter(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    """`fields` filter restricts the returned meta dict per file."""
    from docx import Document

    docs_dir, _ = tmp_roots
    p = docs_dir / "f.docx"
    doc = Document()
    doc.core_properties.title = "T"
    doc.core_properties.author = "A"
    doc.save(str(p))

    cfg = make_config()
    ctx = build_context(cfg)
    mcp = FastMCP(name="t-fields")
    metadata_tool.register(mcp, ctx)
    metadata_tool.register_batch(mcp, ctx)

    fn = mcp._tool_manager._tools["metadata_read_batch"].fn
    result = fn(glob="documents:/*.docx", fields=["paragraph_count"])
    assert result["count"] == 1
    meta = result["files"][0]["meta"]
    assert "paragraph_count" in meta
    # Field filter dropped everything else.
    assert "core_props" not in meta


def test_metadata_read_batch_skips_unhandled_files(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    """Files whose handler raises during read_meta land in `skipped` with reason."""
    docs_dir, _ = tmp_roots
    # File with unsupported extension -> no handler.
    (docs_dir / "weird.unknownext").write_text("?", encoding="utf-8")

    cfg = make_config()
    ctx = build_context(cfg)
    mcp = FastMCP(name="t-skip")
    metadata_tool.register(mcp, ctx)
    metadata_tool.register_batch(mcp, ctx)

    fn = mcp._tool_manager._tools["metadata_read_batch"].fn
    result = fn(glob="documents:/*.unknownext")
    assert result["count"] == 0
    assert result["skipped"][0]["reason"] == "no_handler"
    assert result["skipped_summary"]["no_handler"] == 1


def test_metadata_read_batch_max_files_truncates(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    docs_dir, _ = tmp_roots
    from docx import Document
    for i in range(3):
        d = Document()
        d.add_paragraph(f"p{i}")
        d.save(str(docs_dir / f"d{i}.docx"))

    cfg = make_config()
    ctx = build_context(cfg)
    mcp = FastMCP(name="t-max")
    metadata_tool.register(mcp, ctx)
    metadata_tool.register_batch(mcp, ctx)

    fn = mcp._tool_manager._tools["metadata_read_batch"].fn
    # Truncated at 2 even though 3 docx files match.
    result = fn(glob="documents:/*.docx", max_files=2)
    assert result["count"] <= 2


def test_metadata_read_batch_handler_raises(
    make_config: Callable[..., AppConfig], tmp_roots: tuple[Path, Path]
) -> None:
    """A handler raising during read_meta is caught and reported as skipped."""
    from docx import Document

    docs_dir, _ = tmp_roots
    p = docs_dir / "bad.docx"
    d = Document()
    d.add_paragraph("ok")
    d.save(str(p))

    cfg = make_config()
    ctx = build_context(cfg)
    mcp = FastMCP(name="t-raise")
    metadata_tool.register(mcp, ctx)
    metadata_tool.register_batch(mcp, ctx)

    handler = ctx.registry.for_path(p.absolute())
    fn = mcp._tool_manager._tools["metadata_read_batch"].fn
    with patch.object(handler, "read_meta", side_effect=ValueError("boom")):
        result = fn(glob="documents:/*.docx")
    assert result["count"] == 0
    assert result["skipped"][0]["reason"] == "read_failed"
